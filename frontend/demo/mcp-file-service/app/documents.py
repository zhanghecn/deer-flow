"""Document search and read helpers for the standalone demo MCP service.

This module is intentionally demo-scoped. It gives the 8084 workbench one
reusable document contract without pushing document-specific assumptions into
the generic OpenAgents runtime.
"""

from __future__ import annotations

import base64
import fnmatch
import hashlib
import io
import json
import mimetypes
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable, Literal

from docx import Document as WordDocument
from docx.document import Document as WordDocumentType
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from openpyxl.chart.title import Title as OpenPyxlTitle
from openpyxl.utils import get_column_letter
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader
from xlrd import open_workbook as open_xls_workbook

DocumentKind = Literal[
    "text",
    "image",
    "pdf",
    "ppt",
    "pptx",
    "doc",
    "docx",
    "xls",
    "xlsx",
]
LocatorType = Literal["line", "page", "slide", "sheet", "region"]
EvidenceType = Literal["text", "ocr_text", "table_text", "vision_summary"]
DocumentSearchOutputMode = Literal["content", "files_with_matches", "count"]
NextActionHint = Literal[
    "read_current",
    "read_more",
    "fetch_visual",
    "search_next_page",
]

INLINE_ASSET_LIMIT_BYTES = 200_000
CACHE_SCHEMA_VERSION = 2
DEFAULT_DOCUMENT_SEARCH_HEAD_LIMIT = 250
DOCUMENT_SEARCH_CONTENT_MAX_CONTEXT = 30
DOCUMENT_SEARCH_CONTENT_MAX_MATCHES = 12
DOCUMENT_SEARCH_CONTENT_MAX_BYTES = 8_000
DOCUMENT_READ_MAX_UNITS = 12
DOCUMENT_READ_LOCATOR_MAX_UNITS = 24
DOCUMENT_READ_LOCATOR_DEFAULT_BEFORE = 6
DOCUMENT_READ_LOCATOR_DEFAULT_AFTER = 16
DOCUMENT_READ_TEXT_MAX_BYTES = 700
DOCUMENT_READ_TABLE_MAX_ROWS = 24
DOCUMENT_READ_MAX_BYTES = 12_000
DEFAULT_OCR_LANGUAGES = ("eng", "chi_sim")
PDF_OCR_SCALE = 2.5
IMAGE_EXTENSIONS = {
    ".bmp",
    ".gif",
    ".jpeg",
    ".jpg",
    ".png",
    ".webp",
}
LEGACY_OFFICE_EXTENSIONS = {
    ".doc": "doc",
    ".ppt": "ppt",
    ".xls": "xls",
}
OOXML_OFFICE_EXTENSIONS = {
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
}
GREP_REGEX_TOKENS = ("|", "\\", "[", "]", "(", ")", "{", "}", "^", "$", ".", "*", "+", "?")


@dataclass(frozen=True)
class DocumentAsset:
    """One fetchable or descriptive visual asset surfaced by `document_read`."""

    asset_ref: str
    kind: Literal["image", "document"]
    summary: str
    mime_type: str | None = None
    size_bytes: int | None = None
    width: int | None = None
    height: int | None = None
    data: bytes | None = field(default=None, repr=False)
    extra: dict[str, Any] = field(default_factory=dict)

    def read_block(self, *, locator: str | int) -> dict[str, Any]:
        """Return the compact block shape used by `document_read`."""

        block: dict[str, Any] = {
            "type": "image" if self.kind == "image" else "document",
            "locator": locator,
            "asset_ref": self.asset_ref,
            "summary": self.summary,
            "mime_type": self.mime_type,
            "size_bytes": self.size_bytes,
            "width": self.width,
            "height": self.height,
        }
        if self.extra:
            block["metadata"] = self.extra
        return block

    def fetch_payload(self, *, path: str) -> dict[str, Any]:
        """Inline small binary assets so the demo can verify fetch semantics."""

        payload = {
            "path": path,
            "asset_ref": self.asset_ref,
            "asset_kind": self.kind,
            "summary": self.summary,
            "mime_type": self.mime_type,
            "size_bytes": self.size_bytes,
            "width": self.width,
            "height": self.height,
            "fetchable": self.data is not None,
            "inlined": False,
            "extra": self.extra,
        }
        if self.data is None:
            return payload
        if len(self.data) > INLINE_ASSET_LIMIT_BYTES:
            payload["fetchable"] = False
            payload["warning"] = (
                "asset is larger than the inline demo threshold; "
                "metadata is returned without base64 content"
            )
            return payload
        payload["inlined"] = True
        payload["content_base64"] = base64.b64encode(self.data).decode("ascii")
        return payload


@dataclass(frozen=True)
class DocumentUnit:
    """One addressable unit in a parsed document."""

    locator: str | int
    locator_type: LocatorType
    content_blocks: list[dict[str, Any]]
    search_entries: list[tuple[EvidenceType, str]]
    contains_visual: bool
    preview_text: str


@dataclass(frozen=True)
class ParsedDocument:
    """Structured document view shared by search, read, and asset fetch."""

    path: str
    document_kind: DocumentKind
    locator_type: LocatorType
    units: list[DocumentUnit]
    assets: dict[str, DocumentAsset]
    contains_visual: bool
    ingest_metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(frozen=True)
class OcrText:
    """Normalized OCR text plus the ingest metadata needed by the cache manifest."""

    text: str
    provider: str
    languages: tuple[str, ...]


def _collapse_whitespace(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def _grep_match(
    *,
    pattern: str,
    text: str,
    case_insensitive: bool,
    multiline: bool = False,
) -> tuple[int, str] | None:
    """Return the first grep-style match, using literal search before regex.

    `document_search` intentionally keeps a grep contract. That makes the MCP
    behavior predictable for agents: broad semantic retrieval belongs in a
    future index/reranker tool, while this tool answers "does this exact
    pattern occur in parsed document text/OCR/table evidence?"
    """

    normalized_pattern = pattern.strip()
    if not normalized_pattern:
        return None

    haystack = text.lower() if case_insensitive else text
    needle = normalized_pattern.lower() if case_insensitive else normalized_pattern
    literal_index = haystack.find(needle)
    if literal_index >= 0:
        return literal_index, text[literal_index : literal_index + len(normalized_pattern)]

    if not any(token in pattern for token in GREP_REGEX_TOKENS):
        return None

    try:
        flags = 0
        if case_insensitive:
            flags |= re.IGNORECASE
        if multiline:
            flags |= re.DOTALL | re.MULTILINE
        match = re.search(pattern, text, flags=flags)
    except re.error:
        return None
    if match is None:
        return None
    return match.start(), match.group(0)


def _normalize_document_search_output_mode(value: str) -> DocumentSearchOutputMode:
    """Accept grep-mode aliases while returning one Claude-style mode string."""

    normalized = str(value or "").strip().lower().replace("-", "_").replace(" ", "_")
    if not normalized:
        return "files_with_matches"
    aliases: dict[str, DocumentSearchOutputMode] = {
        "content": "content",
        "matches": "content",
        "match_content": "content",
        "files": "files_with_matches",
        "file": "files_with_matches",
        "filenames": "files_with_matches",
        "paths": "files_with_matches",
        "files_with_matches": "files_with_matches",
        "count": "count",
        "counts": "count",
        "summary": "count",
        "stats": "count",
    }
    return aliases.get(normalized, "files_with_matches")


def _apply_head_window[T](
    items: list[T],
    *,
    head_limit: int | None,
    offset: int,
) -> tuple[list[T], int | None, int | None]:
    """Mirror Claude Code Grep's head/offset contract for all output modes."""

    safe_offset = max(offset, 0)
    if head_limit == 0:
        return items[safe_offset:], None, safe_offset or None
    effective_limit = (
        DEFAULT_DOCUMENT_SEARCH_HEAD_LIMIT
        if head_limit is None
        else max(int(head_limit), 1)
    )
    window = items[safe_offset : safe_offset + effective_limit]
    truncated = len(items) - safe_offset > effective_limit
    return window, effective_limit if truncated else None, safe_offset or None


def _context_line_numbers(
    *,
    total_lines: int,
    line_number: int,
    before: int,
    after: int,
) -> range:
    """Return the one-based grep context line range around a matched line."""

    start = max(line_number - max(before, 0), 1)
    end = min(line_number + max(after, 0), total_lines)
    return range(start, end + 1)


def _build_snippet(*, text: str, match_index: int, window: int = 90) -> str:
    if not text:
        return ""
    if match_index < 0:
        return _collapse_whitespace(text[:window])
    start = max(match_index - window // 3, 0)
    end = min(start + window, len(text))
    return _collapse_whitespace(text[start:end])


def _compact_content_match(item: dict[str, Any]) -> dict[str, Any]:
    """Keep content-mode metadata locator-sized instead of duplicating snippets.

    The human/model-readable grep body already lives in the `content` field.
    Repeating all context lines inside `items/results` can turn a handful of
    broad searches into hundreds of kilobytes of tool output and exhaust the
    next model call's context window.
    """

    read_args = {
        "path": item["path"],
        "locator": item["locator"],
        "before": DOCUMENT_READ_LOCATOR_DEFAULT_BEFORE,
        "after": DOCUMENT_READ_LOCATOR_DEFAULT_AFTER,
    }
    return {
        "path": item["path"],
        "document_kind": item["document_kind"],
        "locator": item["locator"],
        "locator_type": item["locator_type"],
        "line_number": item["line_number"],
        "match_text": item["match_text"],
        "snippet": item["snippet"],
        "evidence_type": item["evidence_type"],
        "contains_visual": item["contains_visual"],
        "next_action_hint": item["next_action_hint"],
        "read_args": read_args,
    }


def _content_match_lines(
    item: dict[str, Any],
    *,
    show_line_numbers: bool,
) -> tuple[list[str], bool]:
    """Render one grep match into bounded source lines plus truncation state."""

    lines: list[str] = []
    truncated = False
    context_lines = item.get("context_lines") or [
        {
            "line_number": item["line_number"],
            "line": item["line"],
            "is_match": True,
        }
    ]
    for context_item in context_lines:
        locator = item["locator"]
        line_number = context_item["line_number"]
        prefix = f"{item['path']}:{item['locator_type']}:{locator}"
        if show_line_numbers:
            prefix = f"{prefix}:{line_number}"
        # Context lines can contain long source URL fields. Keep grep output
        # source-faithful but bounded so one noisy line does not evict every
        # useful neighboring line.
        line_text, line_truncated = _truncate_utf8_text(
            str(context_item["line"]),
            max_bytes=DOCUMENT_READ_TEXT_MAX_BYTES,
        )
        truncated = truncated or line_truncated
        lines.append(f"{prefix}:{line_text}")
    return lines, truncated


def _content_payload_size(payload: dict[str, Any]) -> int:
    """Measure serialized UTF-8 size so tool responses stay model-safe."""

    return len(json.dumps(payload, ensure_ascii=False).encode("utf-8"))


def _truncate_utf8_text(value: str, *, max_bytes: int) -> tuple[str, bool]:
    """Trim text on a UTF-8 boundary so payload caps never corrupt content."""

    encoded = value.encode("utf-8")
    if len(encoded) <= max_bytes:
        return value, False
    marker = "\n[truncated]"
    marker_bytes = marker.encode("utf-8")
    budget = max(0, max_bytes - len(marker_bytes))
    trimmed = encoded[:budget].decode("utf-8", errors="ignore").rstrip()
    return f"{trimmed}{marker}", True


def _compact_read_block(block: dict[str, Any]) -> tuple[dict[str, Any], bool]:
    """Bound one `document_read` block while preserving source locators.

    `document_read` sends raw document text into the next model call. Broad
    "give me all details" requests can otherwise combine many reads into a
    provider-side context or safety failure. The block remains source-faithful
    and advertises truncation so callers can page or narrow explicitly.
    """

    compact = dict(block)
    truncated = False

    text = compact.get("text")
    if isinstance(text, str):
        compact["text"], text_truncated = _truncate_utf8_text(
            text,
            max_bytes=DOCUMENT_READ_TEXT_MAX_BYTES,
        )
        truncated = truncated or text_truncated

    rows = compact.get("rows")
    if isinstance(rows, list) and len(rows) > DOCUMENT_READ_TABLE_MAX_ROWS:
        compact["rows"] = rows[:DOCUMENT_READ_TABLE_MAX_ROWS]
        compact["rows_truncated"] = True
        truncated = True

    summary = compact.get("summary")
    if isinstance(summary, str):
        compact["summary"], summary_truncated = _truncate_utf8_text(
            summary,
            max_bytes=DOCUMENT_READ_TEXT_MAX_BYTES,
        )
        truncated = truncated or summary_truncated

    if truncated:
        compact["truncated"] = True
    return compact, truncated


def _next_action_hint(
    *,
    unit_index: int,
    total_units: int,
    contains_visual: bool,
    evidence_type: EvidenceType,
) -> NextActionHint:
    if evidence_type == "vision_summary" or contains_visual:
        return "fetch_visual"
    if unit_index + 1 < total_units:
        return "read_more"
    return "read_current"


def _find_unit_index_by_locator(
    units: list[DocumentUnit],
    locator: str | int,
) -> int | None:
    """Resolve a source locator to the corresponding parsed unit index.

    Search results expose locators such as text line numbers, PDF pages, and
    slide numbers. Agents can pass that locator back directly instead of
    guessing the zero-based pagination cursor used by broad reads.
    """

    locator_text = str(locator).strip()
    if not locator_text:
        return None
    for index, unit in enumerate(units):
        if str(unit.locator) == locator_text:
            return index
    return None


def _guess_mime_type(name: str) -> str:
    return mimetypes.guess_type(name)[0] or "application/octet-stream"


def _asset_from_bytes(
    *,
    asset_ref: str,
    data: bytes,
    file_name: str,
    summary: str,
    extra: dict[str, Any] | None = None,
) -> DocumentAsset:
    width = None
    height = None
    mime_type = _guess_mime_type(file_name)
    try:
        with Image.open(io.BytesIO(data)) as image:
            width, height = image.size
            mime_type = Image.MIME.get(image.format, mime_type)
    except Exception:
        width = None
        height = None
    return DocumentAsset(
        asset_ref=asset_ref,
        kind="image",
        summary=summary,
        mime_type=mime_type,
        size_bytes=len(data),
        width=width,
        height=height,
        data=data,
        extra=extra or {},
    )


def _asset_from_path(*, asset_ref: str, file_path: Path) -> DocumentAsset:
    data = file_path.read_bytes()
    return _asset_from_bytes(
        asset_ref=asset_ref,
        data=data,
        file_name=file_path.name,
        summary=f"Image file {file_path.name}",
        extra={"source": "standalone_image"},
    )


def _table_rows_to_text(rows: list[list[str]]) -> str:
    lines = []
    for row in rows:
        rendered = " | ".join(cell for cell in row if cell.strip())
        if rendered:
            lines.append(rendered)
    return "\n".join(lines)


def _iter_docx_blocks(document: WordDocumentType) -> list[Paragraph | Table]:
    """Preserve paragraph/table order so region cursors stay meaningful."""

    blocks: list[Paragraph | Table] = []
    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            blocks.append(Paragraph(child, document))
        elif isinstance(child, CT_Tbl):
            blocks.append(Table(child, document))
    return blocks


def _extract_docx_assets(document: WordDocumentType) -> tuple[dict[str, DocumentAsset], list[str]]:
    assets: dict[str, DocumentAsset] = {}
    ordered_refs: list[str] = []
    image_index = 0
    for relation in document.part.rels.values():
        if "image" not in relation.reltype:
            continue
        image_index += 1
        asset_ref = f"media:image:{image_index}"
        image_name = getattr(relation.target_part, "partname", None)
        file_name = Path(str(image_name)).name if image_name else f"image-{image_index}.bin"
        assets[asset_ref] = _asset_from_bytes(
            asset_ref=asset_ref,
            data=relation.target_part.blob,
            file_name=file_name,
            summary=f"Embedded image {image_index}",
            extra={"source": "docx"},
        )
        ordered_refs.append(asset_ref)
    return assets, ordered_refs


def _chart_summary_text(title: Any, fallback: str) -> str:
    if title is None:
        return fallback
    text = ""
    if hasattr(title, "text_frame"):
        text = getattr(title.text_frame, "text", "") or ""
    elif isinstance(title, OpenPyxlTitle):
        rich = getattr(getattr(title, "tx", None), "rich", None)
        if rich and getattr(rich, "p", None):
            parts: list[str] = []
            for paragraph in rich.p:
                for run in getattr(paragraph, "r", []) or []:
                    if getattr(run, "t", None):
                        parts.append(str(run.t))
                if getattr(paragraph, "endParaRPr", None):
                    continue
            text = " ".join(parts)
    return _collapse_whitespace(text) or fallback


def _extract_sheet_rows(sheet: Any) -> tuple[list[list[str]], str]:
    rows: list[list[str]] = []
    max_column = 0
    for row in sheet.iter_rows():
        rendered = ["" if cell.value is None else _collapse_whitespace(str(cell.value)) for cell in row]
        if any(rendered):
            rows.append(rendered)
            max_column = max(max_column, len(rendered))
    if not rows:
        return [], f"{sheet.title}!A1"
    end_ref = f"{get_column_letter(max_column)}{len(rows)}"
    return rows, f"{sheet.title}!A1:{end_ref}"


def _extract_xls_rows(sheet: Any) -> tuple[list[list[str]], str]:
    rows: list[list[str]] = []
    max_column = 0
    for row_index in range(sheet.nrows):
        rendered = [
            _collapse_whitespace(str(sheet.cell_value(row_index, column_index)))
            for column_index in range(sheet.ncols)
        ]
        if any(rendered):
            rows.append(rendered)
            max_column = max(max_column, len(rendered))
    if not rows:
        return [], f"{sheet.name}!A1"
    end_ref = f"{get_column_letter(max_column)}{len(rows)}"
    return rows, f"{sheet.name}!A1:{end_ref}"


def _cache_safe_name(value: str) -> str:
    return re.sub(r"[^a-zA-Z0-9._-]+", "_", value).strip("._") or "asset"


def _asset_extension(asset: DocumentAsset) -> str:
    guessed = mimetypes.guess_extension(asset.mime_type or "")
    if guessed:
        return guessed
    return ".bin"


def _locator_heading(locator_type: LocatorType, locator: str | int) -> str:
    labels = {
        "line": "Line",
        "page": "Page",
        "slide": "Slide",
        "sheet": "Sheet",
        "region": "Region",
    }
    return f"{labels.get(locator_type, locator_type.title())} {locator}"


class DocumentTooling:
    """Parse uploaded files into the demo's richer document contract."""

    def __init__(
        self,
        *,
        root: Path,
        describe_access: Callable[[Path], Any],
        cache_root: Path,
        ocr_languages: tuple[str, ...] = DEFAULT_OCR_LANGUAGES,
        tesseract_binary: str = "tesseract",
    ) -> None:
        self.root = root
        self.describe_access = describe_access
        # Cache files are MCP-internal implementation details. Keep them out of
        # the uploaded corpus tree so agents cannot confuse cached markdown with
        # the external knowledge-base contract.
        self.cache_root = cache_root.resolve()
        self.cache_root.mkdir(parents=True, exist_ok=True)
        self.ocr_languages = tuple(language for language in ocr_languages if language)
        self.tesseract_binary = tesseract_binary.strip() or "tesseract"
        self.ocr_available = shutil.which(self.tesseract_binary) is not None

    def parse_document(self, file_path: Path) -> ParsedDocument:
        """Return the parsed document from the MCP-managed cache."""

        return self._load_cached_document(file_path)

    def prime_cache(self) -> None:
        """Warm cached manifests/markdown/assets for every supported document."""

        for file_path in sorted(self.root.rglob("*"), key=lambda item: item.as_posix().lower()):
            if not file_path.is_file():
                continue
            try:
                self.prepare_cached_document(file_path)
            except ValueError:
                continue

    def prepare_cached_document(self, file_path: Path) -> ParsedDocument:
        """Build or refresh one cached document package when the source changed."""

        file_path = file_path.resolve()
        if not file_path.is_file():
            raise ValueError(f"path is not a file: {file_path}")
        source_metadata = self._source_metadata(file_path)
        cache_dir = self._cache_dir(file_path)
        manifest_path = cache_dir / "manifest.json"
        if manifest_path.exists():
            try:
                manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
                if self._manifest_is_fresh(manifest=manifest, source_metadata=source_metadata):
                    return self._document_from_manifest(manifest=manifest, cache_dir=cache_dir)
            except Exception:
                shutil.rmtree(cache_dir, ignore_errors=True)

        document = self._parse_source_document(file_path)
        manifest = self._write_cached_document(
            file_path=file_path,
            document=document,
            source_metadata=source_metadata,
            cache_dir=cache_dir,
        )
        return self._document_from_manifest(manifest=manifest, cache_dir=cache_dir)

    def invalidate_cached_document(self, file_path: Path) -> None:
        """Drop the cached package for one source file after delete/replace."""

        shutil.rmtree(self._cache_dir(file_path.resolve()), ignore_errors=True)

    def _parse_source_document(self, file_path: Path) -> ParsedDocument:
        relative_path = file_path.relative_to(self.root).as_posix()
        suffix = file_path.suffix.lower()
        access = self.describe_access(file_path)
        if suffix in IMAGE_EXTENSIONS:
            return self._parse_image_document(file_path=file_path, path=relative_path)
        if suffix == ".pdf":
            return self._parse_pdf_document(file_path=file_path, path=relative_path)
        if suffix == ".pptx":
            return self._parse_pptx_document(file_path=file_path, path=relative_path)
        if suffix == ".docx":
            return self._parse_docx_document(file_path=file_path, path=relative_path)
        if suffix == ".xlsx":
            return self._parse_xlsx_document(file_path=file_path, path=relative_path)
        if suffix == ".xls":
            return self._parse_xls_document(file_path=file_path, path=relative_path)
        if suffix in LEGACY_OFFICE_EXTENSIONS:
            return self._parse_legacy_document(file_path=file_path, path=relative_path)
        if access.text_readable:
            return self._parse_text_document(file_path=file_path, path=relative_path)
        raise ValueError(
            f"document tools do not support '{relative_path}' ({access.mime_type})"
        )

    def _source_metadata(self, file_path: Path) -> dict[str, Any]:
        stat = file_path.stat()
        digest = hashlib.sha256(file_path.read_bytes()).hexdigest()
        return {
            "path": file_path.relative_to(self.root).as_posix(),
            "size_bytes": stat.st_size,
            "mtime_ns": stat.st_mtime_ns,
            "sha256": digest,
        }

    def _cache_dir(self, file_path: Path) -> Path:
        relative = file_path.relative_to(self.root)
        return self.cache_root.joinpath(*relative.parts)

    @staticmethod
    def _manifest_is_fresh(*, manifest: dict[str, Any], source_metadata: dict[str, Any]) -> bool:
        source = manifest.get("source")
        if not isinstance(source, dict):
            return False
        return (
            int(manifest.get("cache_schema_version", 0)) == CACHE_SCHEMA_VERSION
            and source.get("path") == source_metadata["path"]
            and int(source.get("size_bytes", -1)) == int(source_metadata["size_bytes"])
            and int(source.get("mtime_ns", -1)) == int(source_metadata["mtime_ns"])
            and source.get("sha256") == source_metadata["sha256"]
        )

    def _write_cached_document(
        self,
        *,
        file_path: Path,
        document: ParsedDocument,
        source_metadata: dict[str, Any],
        cache_dir: Path,
    ) -> dict[str, Any]:
        if cache_dir.exists():
            shutil.rmtree(cache_dir)
        assets_dir = cache_dir / "assets"
        assets_dir.mkdir(parents=True, exist_ok=True)

        serialized_assets: dict[str, Any] = {}
        for asset_ref, asset in document.assets.items():
            record: dict[str, Any] = {
                "asset_ref": asset.asset_ref,
                "kind": asset.kind,
                "summary": asset.summary,
                "mime_type": asset.mime_type,
                "size_bytes": asset.size_bytes,
                "width": asset.width,
                "height": asset.height,
                "extra": asset.extra,
            }
            if asset.data is not None:
                asset_name = f"{_cache_safe_name(asset_ref)}{_asset_extension(asset)}"
                asset_path = assets_dir / asset_name
                asset_path.write_bytes(asset.data)
                record["cache_file"] = f"assets/{asset_name}"
            serialized_assets[asset_ref] = record

        manifest = {
            "cache_schema_version": CACHE_SCHEMA_VERSION,
            "source": source_metadata,
            "document": {
                "path": document.path,
                "document_kind": document.document_kind,
                "locator_type": document.locator_type,
                "contains_visual": document.contains_visual,
            },
            # Keep OCR state explicit so the external MCP remains the single
            # document-ingest owner instead of pushing hidden behavior into the runtime.
            "ingest": document.ingest_metadata,
            "units": [
                {
                    "locator": unit.locator,
                    "locator_type": unit.locator_type,
                    "content_blocks": unit.content_blocks,
                    "search_entries": [list(entry) for entry in unit.search_entries],
                    "contains_visual": unit.contains_visual,
                    "preview_text": unit.preview_text,
                }
                for unit in document.units
            ],
            "assets": serialized_assets,
        }
        canonical_markdown, canonical_source = self._canonical_markdown(
            file_path=file_path,
            document=document,
            serialized_assets=serialized_assets,
        )
        manifest["ingest"]["canonical_source"] = canonical_source
        (cache_dir / "canonical.md").write_text(canonical_markdown, encoding="utf-8")
        (cache_dir / "manifest.json").write_text(
            json.dumps(manifest, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return manifest

    def _document_from_manifest(
        self,
        *,
        manifest: dict[str, Any],
        cache_dir: Path,
    ) -> ParsedDocument:
        document_payload = manifest["document"]
        assets: dict[str, DocumentAsset] = {}
        for asset_ref, record in manifest.get("assets", {}).items():
            data = None
            cache_file = record.get("cache_file")
            if isinstance(cache_file, str) and cache_file.strip():
                asset_path = cache_dir / cache_file
                if asset_path.exists():
                    data = asset_path.read_bytes()
            assets[asset_ref] = DocumentAsset(
                asset_ref=record["asset_ref"],
                kind=record["kind"],
                summary=record["summary"],
                mime_type=record.get("mime_type"),
                size_bytes=record.get("size_bytes"),
                width=record.get("width"),
                height=record.get("height"),
                data=data,
                extra=record.get("extra") or {},
            )

        units = [
            DocumentUnit(
                locator=unit["locator"],
                locator_type=unit["locator_type"],
                content_blocks=unit.get("content_blocks") or [],
                search_entries=[
                    (entry[0], entry[1])
                    for entry in unit.get("search_entries") or []
                    if isinstance(entry, list) and len(entry) == 2
                ],
                contains_visual=bool(unit.get("contains_visual")),
                preview_text=str(unit.get("preview_text") or ""),
            )
            for unit in manifest.get("units") or []
            if isinstance(unit, dict)
        ]
        return ParsedDocument(
            path=document_payload["path"],
            document_kind=document_payload["document_kind"],
            locator_type=document_payload["locator_type"],
            units=units,
            assets=assets,
            contains_visual=bool(document_payload.get("contains_visual")),
            ingest_metadata=manifest.get("ingest") or {},
        )

    def _load_cached_document(self, file_path: Path) -> ParsedDocument:
        """Read a document package from cache, rebuilding it on demand."""

        return self.prepare_cached_document(file_path)

    def _default_ingest_metadata(
        self,
        *,
        ocr_attempted: bool,
        ocr_provider: str | None = None,
        canonical_source: str = "parsed_document",
    ) -> dict[str, Any]:
        """Keep OCR/canonical provenance explicit inside the MCP-owned cache."""

        if ocr_attempted:
            ocr_status = "complete" if self.ocr_available else "not_available"
        else:
            ocr_status = "not_needed"
        return {
            "ocr_status": ocr_status,
            "ocr_languages": list(self.ocr_languages) if ocr_attempted else [],
            "ocr_provider": ocr_provider if ocr_attempted else None,
            "canonical_source": canonical_source,
        }

    def _run_tesseract_path(self, file_path: Path) -> OcrText | None:
        """Run OCR through the local Tesseract binary so the demo matches prod-like MCP ownership.

        The OCR responsibility intentionally stays inside the external MCP demo.
        The generic runtime only sees the resulting document contract.
        """

        if not self.ocr_available:
            return None
        command = [
            self.tesseract_binary,
            str(file_path),
            "stdout",
            "-l",
            "+".join(self.ocr_languages),
            "--psm",
            "6",
        ]
        try:
            completed = subprocess.run(
                command,
                check=False,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="ignore",
            )
        except OSError:
            return None
        if completed.returncode != 0:
            return None
        text = _collapse_whitespace(completed.stdout)
        if not text:
            return None
        return OcrText(
            text=text,
            provider="tesseract",
            languages=self.ocr_languages,
        )

    def _run_tesseract_bytes(self, *, data: bytes, suffix: str) -> OcrText | None:
        """OCR in-memory assets by materializing a short-lived temp file for Tesseract."""

        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as handle:
            handle.write(data)
            temp_path = Path(handle.name)
        try:
            return self._run_tesseract_path(temp_path)
        finally:
            temp_path.unlink(missing_ok=True)

    def _ocr_asset_bytes(
        self,
        *,
        data: bytes | None,
        asset_ref: str,
        fallback_suffix: str = ".png",
    ) -> OcrText | None:
        """Run OCR against one extracted image asset when the cache needs semantic text."""

        if data is None:
            return None
        suffix = Path(asset_ref.replace(":", "_")).suffix or fallback_suffix
        return self._run_tesseract_bytes(data=data, suffix=suffix)

    def _render_pdf_page_image(self, *, file_path: Path, page_index: int) -> bytes | None:
        """Render one PDF page to PNG so scanned pages become OCR-able and fetchable."""

        try:
            import fitz
        except Exception:
            return None

        document = fitz.open(str(file_path))
        try:
            page = document.load_page(page_index)
            matrix = fitz.Matrix(PDF_OCR_SCALE, PDF_OCR_SCALE)
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            return pixmap.tobytes("png")
        finally:
            document.close()

    def _canonical_markdown(
        self,
        *,
        file_path: Path,
        document: ParsedDocument,
        serialized_assets: dict[str, Any],
    ) -> tuple[str, str]:
        """Prefer MarkItDown output when available so the cache matches KB-style packages."""

        try:
            from markitdown import MarkItDown
        except Exception:
            return (
                self._render_canonical_markdown(
                    document=document,
                    serialized_assets=serialized_assets,
                ),
                "parsed_document",
            )

        try:
            markdown = str(MarkItDown().convert(str(file_path)).text_content or "").strip()
        except Exception:
            markdown = ""
        if not markdown:
            return (
                self._render_canonical_markdown(
                    document=document,
                    serialized_assets=serialized_assets,
                ),
                "parsed_document",
            )
        return f"# {document.path}\n\n{markdown}\n", "markitdown"

    def _render_canonical_markdown(
        self,
        *,
        document: ParsedDocument,
        serialized_assets: dict[str, Any],
    ) -> str:
        lines = [
            f"# {document.path}",
            "",
            f"- document_kind: {document.document_kind}",
            f"- locator_type: {document.locator_type}",
            f"- contains_visual: {str(document.contains_visual).lower()}",
            "",
        ]
        for unit in document.units:
            lines.extend(
                [
                    f"## {_locator_heading(unit.locator_type, unit.locator)}",
                    "",
                ]
            )
            if not unit.content_blocks:
                lines.extend(["_No content blocks_", ""])
                continue
            for block in unit.content_blocks:
                block_type = block.get("type")
                if block_type == "text":
                    lines.extend([str(block.get("text", "")), ""])
                    continue
                if block_type == "table":
                    lines.extend(
                        [
                            "```table",
                            str(block.get("text", "")),
                            "```",
                            "",
                        ]
                    )
                    continue
                if block_type == "image":
                    asset_ref = str(block.get("asset_ref", ""))
                    record = serialized_assets.get(asset_ref) or {}
                    cache_file = record.get("cache_file")
                    if isinstance(cache_file, str) and cache_file.strip():
                        lines.append(f"![{block.get('summary', asset_ref)}]({cache_file})")
                    else:
                        lines.append(f"- image: {block.get('summary', asset_ref)}")
                    lines.append("")
                    continue
                lines.extend([f"> {block.get('summary', '')}".rstrip(), ""])
        return "\n".join(lines).strip() + "\n"

    def describe_document(self, file_path: Path) -> dict[str, Any]:
        """Return lightweight document metadata for browse/list surfaces.

        `document_list` should be able to enumerate the MCP-owned knowledge tree
        without forcing the agent to guess whether a file is searchable. Keep
        unsupported files explicit instead of silently hiding them.
        """

        relative_path = file_path.relative_to(self.root).as_posix()
        access = self.describe_access(file_path)
        try:
            document = self.parse_document(file_path)
        except ValueError:
            return {
                "path": relative_path,
                "document_kind": None,
                "locator_type": None,
                "contains_visual": False,
                "supported": False,
                "mime_type": access.mime_type,
            }
        return {
            "path": document.path,
            "document_kind": document.document_kind,
            "locator_type": document.locator_type,
            "contains_visual": document.contains_visual,
            "supported": True,
            "mime_type": access.mime_type,
        }

    def search(
        self,
        *,
        files: list[Path],
        pattern: str,
        output_mode: str = "files_with_matches",
        glob: str | None = None,
        cursor: int = 0,
        limit: int = 10,
        context_before: int | None = None,
        context_after: int | None = None,
        context: int | None = None,
        show_line_numbers: bool = True,
        case_insensitive: bool = False,
        head_limit: int | None = None,
        offset: int = 0,
        multiline: bool = False,
    ) -> dict[str, Any]:
        search_pattern = pattern.strip()
        if not search_pattern:
            raise ValueError("pattern is required")
        resolved_mode = _normalize_document_search_output_mode(output_mode)
        effective_offset = offset if offset > 0 else max(cursor, 0)
        if head_limit is None and limit != 10:
            # `cursor`/`limit` are legacy demo arguments. Keep explicit old
            # callers working without weakening the Claude-style head_limit API.
            head_limit = limit
        requested_before = max(context_before if context_before is not None else context or 0, 0)
        requested_after = max(context_after if context_after is not None else context or 0, 0)
        if resolved_mode == "content":
            # Content mode feeds directly back into the next model turn. Cap
            # context like Claude-style grep output so a broad "give me all"
            # request paginates instead of exhausting the model window.
            before = min(requested_before, DOCUMENT_SEARCH_CONTENT_MAX_CONTEXT)
            after = min(requested_after, DOCUMENT_SEARCH_CONTENT_MAX_CONTEXT)
        else:
            before = requested_before
            after = requested_after

        matches: list[dict[str, Any]] = []
        skipped_files: list[str] = []
        for file_path in files:
            relative_path = file_path.relative_to(self.root).as_posix()
            if glob and not fnmatch.fnmatch(relative_path, glob):
                continue
            try:
                document = self.parse_document(file_path)
            except ValueError:
                skipped_files.append(relative_path)
                continue

            for unit_index, unit in enumerate(document.units):
                for evidence_type, evidence_text in unit.search_entries:
                    if multiline:
                        match = _grep_match(
                            pattern=search_pattern,
                            text=evidence_text,
                            case_insensitive=case_insensitive,
                            multiline=True,
                        )
                        if match is None:
                            continue
                        match_index, match_text = match
                        line_number = evidence_text[:match_index].count("\n") + 1
                        lines = evidence_text.splitlines() or [evidence_text]
                        line = lines[line_number - 1] if line_number <= len(lines) else ""
                        matches.append(
                            {
                                "path": document.path,
                                "document_kind": document.document_kind,
                                "locator": unit.locator,
                                "locator_type": unit.locator_type,
                                "line_number": line_number,
                                "line": line,
                                "match_text": match_text,
                                "snippet": _build_snippet(
                                    text=evidence_text,
                                    match_index=match_index,
                                ),
                                "context_lines": [
                                    {
                                        "line_number": context_line_number,
                                        "line": lines[context_line_number - 1],
                                        "is_match": context_line_number == line_number,
                                    }
                                    for context_line_number in _context_line_numbers(
                                        total_lines=len(lines),
                                        line_number=line_number,
                                        before=before,
                                        after=after,
                                    )
                                ],
                                "evidence_type": evidence_type,
                                "contains_visual": unit.contains_visual,
                                "next_action_hint": _next_action_hint(
                                    unit_index=unit_index,
                                    total_units=len(document.units),
                                    contains_visual=unit.contains_visual,
                                    evidence_type=evidence_type,
                                ),
                            }
                        )
                        continue

                    lines = evidence_text.splitlines() or [evidence_text]
                    for line_number, line in enumerate(lines, start=1):
                        match = _grep_match(
                            pattern=search_pattern,
                            text=line,
                            case_insensitive=case_insensitive,
                        )
                        if match is None:
                            continue
                        match_index, match_text = match
                        matches.append(
                            {
                                "path": document.path,
                                "document_kind": document.document_kind,
                                "locator": unit.locator,
                                "locator_type": unit.locator_type,
                                "line_number": line_number,
                                "line": line,
                                "match_text": match_text,
                                "snippet": _build_snippet(
                                    text=line,
                                    match_index=match_index,
                                ),
                                "context_lines": [
                                    {
                                        "line_number": context_line_number,
                                        "line": lines[context_line_number - 1],
                                        "is_match": context_line_number == line_number,
                                    }
                                    for context_line_number in _context_line_numbers(
                                        total_lines=len(lines),
                                        line_number=line_number,
                                        before=before,
                                        after=after,
                                    )
                                ],
                                "evidence_type": evidence_type,
                                "contains_visual": unit.contains_visual,
                                "next_action_hint": _next_action_hint(
                                    unit_index=unit_index,
                                    total_units=len(document.units),
                                    contains_visual=unit.contains_visual,
                                    evidence_type=evidence_type,
                                ),
                            }
                        )

        matches.sort(
            key=lambda item: (
                item["path"],
                str(item["locator"]),
                str(item["evidence_type"]),
                int(item["line_number"]),
            )
        )
        metadata: dict[str, Any] = {
            "pattern": search_pattern,
            "mode": resolved_mode,
            "output_mode": resolved_mode,
            "requested_output_mode": output_mode,
            "case_insensitive": case_insensitive,
            "multiline": multiline,
            "skipped_files": skipped_files,
        }
        if resolved_mode == "content":
            metadata["requested_context_before"] = requested_before
            metadata["requested_context_after"] = requested_after
            metadata["applied_context_before"] = before
            metadata["applied_context_after"] = after
            if before != requested_before or after != requested_after:
                metadata["context_truncated"] = True
        if glob:
            metadata["glob"] = glob

        if resolved_mode == "files_with_matches":
            filenames = sorted({item["path"] for item in matches})
            window, applied_limit, applied_offset = _apply_head_window(
                filenames,
                head_limit=head_limit,
                offset=effective_offset,
            )
            payload = {
                **metadata,
                "filenames": window,
                "numFiles": len(filenames),
                "total": len(filenames),
                "items": window,
                "results": window,
            }
        elif resolved_mode == "count":
            file_counts: dict[str, int] = {}
            for item in matches:
                file_counts[item["path"]] = file_counts.get(item["path"], 0) + 1
            count_rows = [
                {"path": path, "match_count": count}
                for path, count in sorted(file_counts.items())
            ]
            window, applied_limit, applied_offset = _apply_head_window(
                count_rows,
                head_limit=head_limit,
                offset=effective_offset,
            )
            content_lines = [
                f"{item['path']}:{item['match_count']}"
                for item in window
            ]
            payload = {
                **metadata,
                "filenames": [item["path"] for item in count_rows],
                "numFiles": len(count_rows),
                "numMatches": len(matches),
                "content": "\n".join(content_lines),
                "items": window,
                "results": window,
                "total": len(count_rows),
                "total_matches": len(matches),
            }
        else:
            if head_limit == 0:
                content_head_limit = DOCUMENT_SEARCH_CONTENT_MAX_MATCHES
            elif head_limit is None:
                content_head_limit = min(
                    DEFAULT_DOCUMENT_SEARCH_HEAD_LIMIT,
                    DOCUMENT_SEARCH_CONTENT_MAX_MATCHES,
                )
            else:
                content_head_limit = min(
                    max(int(head_limit), 1),
                    DOCUMENT_SEARCH_CONTENT_MAX_MATCHES,
                )
            window, applied_limit, applied_offset = _apply_head_window(
                matches,
                head_limit=content_head_limit,
                offset=effective_offset,
            )
            match_line_groups: list[list[str]] = []
            content_was_truncated = False
            for item in window:
                lines, lines_truncated = _content_match_lines(
                    item,
                    show_line_numbers=show_line_numbers,
                )
                match_line_groups.append(lines)
                content_was_truncated = content_was_truncated or lines_truncated
            compact_window = [_compact_content_match(item) for item in window]
            content_lines = [line for group in match_line_groups for line in group]
            payload = {
                **metadata,
                "filenames": sorted({item["path"] for item in matches}),
                "numFiles": len({item["path"] for item in matches}),
                "numLines": len(content_lines),
                "content": "\n".join(content_lines),
                "items": compact_window,
                "results": compact_window,
                "total": len(matches),
            }
            while len(compact_window) > 1 and _content_payload_size(payload) > DOCUMENT_SEARCH_CONTENT_MAX_BYTES:
                compact_window.pop()
                match_line_groups.pop()
                content_lines = [line for group in match_line_groups for line in group]
                payload["content"] = "\n".join(content_lines)
                payload["numLines"] = len(content_lines)
                payload["items"] = compact_window
                payload["results"] = compact_window
                content_was_truncated = True
            if compact_window and _content_payload_size(payload) > DOCUMENT_SEARCH_CONTENT_MAX_BYTES:
                first_item = window[0]
                line_group, _ = _content_match_lines(
                    {
                        **first_item,
                        "context_lines": [
                            {
                                "line_number": first_item["line_number"],
                                "line": first_item["line"],
                                "is_match": True,
                            }
                        ],
                    },
                    show_line_numbers=show_line_numbers,
                )
                content_lines = line_group
                payload["content"] = "\n".join(content_lines)
                payload["numLines"] = len(content_lines)
                payload["items"] = compact_window[:1]
                payload["results"] = compact_window[:1]
                compact_window = compact_window[:1]
                content_was_truncated = True
            if _content_payload_size(payload) > DOCUMENT_SEARCH_CONTENT_MAX_BYTES:
                payload["content"] = ""
                payload["numLines"] = 0
                content_was_truncated = True
            if content_was_truncated or len(window) < len(matches) or len(content_lines) < sum(
                len(
                    item.get("context_lines")
                    or [
                        {
                            "line_number": item["line_number"],
                            "line": item["line"],
                            "is_match": True,
                        }
                    ]
                )
                for item in window
            ):
                payload["content_truncated"] = True
                payload["max_content_bytes"] = DOCUMENT_SEARCH_CONTENT_MAX_BYTES

        if applied_limit is not None:
            payload["appliedLimit"] = applied_limit
        if applied_offset is not None:
            payload["appliedOffset"] = applied_offset
        payload["has_more"] = payload["total"] > effective_offset + len(payload.get("items", []))
        payload["next_offset"] = (
            effective_offset + len(payload.get("items", []))
            if payload["has_more"]
            else None
        )
        return payload

    def read(
        self,
        *,
        file_path: Path,
        cursor: int,
        limit: int,
        locator: str | int | None = None,
        before: int | None = None,
        after: int | None = None,
    ) -> dict[str, Any]:
        document = self.parse_document(file_path)
        safe_cursor = max(cursor, 0)
        requested_limit = max(limit, 1)
        total_units = len(document.units)
        requested_locator = None if locator is None else str(locator).strip()
        locator_index: int | None = None
        locator_window = False
        requested_before = max(before or 0, 0)
        requested_after = max(after or 0, 0)
        applied_before = 0
        applied_after = 0

        if requested_locator:
            locator_index = _find_unit_index_by_locator(document.units, requested_locator)
            if locator_index is None:
                raise ValueError(
                    f"locator '{requested_locator}' does not exist for {document.path}"
                )
            locator_window = True
            if before is None and after is None:
                requested_before = DOCUMENT_READ_LOCATOR_DEFAULT_BEFORE
                requested_after = DOCUMENT_READ_LOCATOR_DEFAULT_AFTER
            elif after is None:
                requested_after = max(requested_limit - 1, 0)
            requested_span = requested_before + 1 + requested_after
            safe_span = min(max(requested_span, 1), DOCUMENT_READ_LOCATOR_MAX_UNITS)
            applied_before = min(requested_before, safe_span - 1, locator_index)
            remaining_after_budget = max(safe_span - applied_before - 1, 0)
            applied_after = min(
                requested_after,
                remaining_after_budget,
                max(total_units - locator_index - 1, 0),
            )
            safe_cursor = max(locator_index - applied_before, 0)
            safe_limit = applied_before + 1 + applied_after
        else:
            safe_limit = min(requested_limit, DOCUMENT_READ_MAX_UNITS)
        if total_units == 0:
            return {
                "path": document.path,
                "document_kind": document.document_kind,
                "locator_type": document.locator_type,
                "cursor": safe_cursor,
                "limit": safe_limit,
                "requested_limit": requested_limit,
                "has_more": False,
                "next_cursor": None,
                "total_units": 0,
                "returned_units": 0,
                "contains_visual": document.contains_visual,
                "content_blocks": [],
            }

        selected_units = document.units[safe_cursor : safe_cursor + safe_limit]
        content_blocks: list[dict[str, Any]] = []
        content_truncated = (
            safe_limit != requested_limit
            if not locator_window
            else safe_limit != requested_before + 1 + requested_after
        )
        for unit in selected_units:
            for block in unit.content_blocks:
                compact_block, block_truncated = _compact_read_block(block)
                content_blocks.append(compact_block)
                content_truncated = content_truncated or block_truncated

        next_cursor = safe_cursor + len(selected_units)
        payload = {
            "path": document.path,
            "document_kind": document.document_kind,
            "locator_type": document.locator_type,
            "cursor": safe_cursor,
            "limit": safe_limit,
            "requested_limit": requested_limit,
            "has_more": next_cursor < total_units,
            "next_cursor": next_cursor if next_cursor < total_units else None,
            "total_units": total_units,
            "returned_units": len(selected_units),
            "contains_visual": any(unit.contains_visual for unit in selected_units),
            "content_blocks": content_blocks,
        }
        if locator_window:
            payload["locator"] = requested_locator
            payload["locator_window"] = True
            payload["requested_before"] = requested_before
            payload["requested_after"] = requested_after
            payload["applied_before"] = applied_before
            payload["applied_after"] = applied_after
            payload["matched_cursor"] = locator_index
        if (not locator_window and safe_limit != requested_limit) or (
            locator_window and safe_limit != requested_before + 1 + requested_after
        ):
            payload["limit_truncated"] = True

        # Most payloads are bounded by unit and block caps. This final guard
        # handles unusual mixed media/table documents with large metadata.
        while (
            len(payload["content_blocks"]) > 1
            and _content_payload_size(payload) > DOCUMENT_READ_MAX_BYTES
        ):
            payload["content_blocks"].pop()
            payload["content_truncated"] = True

        if content_truncated or payload.get("content_truncated"):
            payload["content_truncated"] = True
            payload["max_content_bytes"] = DOCUMENT_READ_MAX_BYTES
        return payload

    def fetch_asset(self, *, file_path: Path, asset_ref: str) -> dict[str, Any]:
        if not asset_ref.strip():
            raise ValueError("asset_ref is required")
        document = self.parse_document(file_path)
        asset = document.assets.get(asset_ref)
        if asset is None:
            raise FileNotFoundError(
                f"asset_ref '{asset_ref}' does not exist for {document.path}"
            )
        return asset.fetch_payload(path=document.path)

    def _parse_text_document(self, *, file_path: Path, path: str) -> ParsedDocument:
        lines = file_path.read_text(encoding="utf-8", errors="ignore").splitlines()
        units: list[DocumentUnit] = []
        if not lines:
            units.append(
                DocumentUnit(
                    locator=1,
                    locator_type="line",
                    content_blocks=[
                        {
                            "type": "text",
                            "locator": 1,
                            "text": "",
                        }
                    ],
                    search_entries=[],
                    contains_visual=False,
                    preview_text="",
                )
            )
        for line_number, line in enumerate(lines, start=1):
            units.append(
                DocumentUnit(
                    locator=line_number,
                    locator_type="line",
                    content_blocks=[
                        {
                            "type": "text",
                            "locator": line_number,
                            "text": line,
                        }
                    ],
                    search_entries=[("text", line)] if line.strip() else [],
                    contains_visual=False,
                    preview_text=line,
                )
            )
        return ParsedDocument(
            path=path,
            document_kind="text",
            locator_type="line",
            units=units,
            assets={},
            contains_visual=False,
            ingest_metadata=self._default_ingest_metadata(ocr_attempted=False),
        )

    def _parse_image_document(self, *, file_path: Path, path: str) -> ParsedDocument:
        asset = _asset_from_path(asset_ref="page:1:image:1", file_path=file_path)
        summary = f"Image file {file_path.name} ({asset.width or '?'}x{asset.height or '?'})"
        ocr_text = self._run_tesseract_path(file_path)
        content_blocks: list[dict[str, Any]] = []
        search_entries: list[tuple[EvidenceType, str]] = []
        if ocr_text is not None:
            content_blocks.append(
                {
                    "type": "text",
                    "locator": 1,
                    "text": ocr_text.text,
                    "text_source": "ocr",
                }
            )
            search_entries.append(("ocr_text", ocr_text.text))
        unit = DocumentUnit(
            locator=1,
            locator_type="page",
            content_blocks=[
                *content_blocks,
                asset.read_block(locator=1),
                {
                    "type": "document",
                    "locator": 1,
                    "summary": summary,
                },
            ],
            search_entries=[
                *search_entries,
                ("vision_summary", f"{file_path.stem} {summary}"),
            ],
            contains_visual=True,
            preview_text=ocr_text.text if ocr_text is not None else summary,
        )
        return ParsedDocument(
            path=path,
            document_kind="image",
            locator_type="page",
            units=[unit],
            assets={asset.asset_ref: asset},
            contains_visual=True,
            ingest_metadata=self._default_ingest_metadata(
                ocr_attempted=True,
                ocr_provider=ocr_text.provider if ocr_text is not None else "tesseract",
            ),
        )

    def _parse_pdf_document(self, *, file_path: Path, path: str) -> ParsedDocument:
        reader = PdfReader(str(file_path))
        units: list[DocumentUnit] = []
        assets: dict[str, DocumentAsset] = {}
        ocr_attempted = False
        for page_index, page in enumerate(reader.pages, start=1):
            text = _collapse_whitespace(page.extract_text() or "")
            content_blocks: list[dict[str, Any]] = []
            search_entries: list[tuple[EvidenceType, str]] = []
            asset_refs: list[str] = []

            if text:
                content_blocks.append(
                    {
                        "type": "text",
                        "locator": page_index,
                        "text": text,
                    }
                )
                search_entries.append(("text", text))

            for image_index, image_file in enumerate(page.images, start=1):
                asset_ref = f"page:{page_index}:image:{image_index}"
                asset = _asset_from_bytes(
                    asset_ref=asset_ref,
                    data=image_file.data,
                    file_name=image_file.name,
                    summary=f"Embedded image on page {page_index}",
                    extra={"page": page_index, "source": "pdf"},
                )
                assets[asset_ref] = asset
                asset_refs.append(asset_ref)
                content_blocks.append(asset.read_block(locator=page_index))

            ocr_text: OcrText | None = None
            if not text:
                ocr_attempted = True
                rendered_page = self._render_pdf_page_image(
                    file_path=file_path,
                    page_index=page_index - 1,
                )
                if rendered_page:
                    if not asset_refs:
                        asset_ref = f"page:{page_index}:render:1"
                        asset = _asset_from_bytes(
                            asset_ref=asset_ref,
                            data=rendered_page,
                            file_name=f"page-{page_index}.png",
                            summary=f"Rendered page {page_index}",
                            extra={"page": page_index, "source": "pdf_render"},
                        )
                        assets[asset_ref] = asset
                        asset_refs.append(asset_ref)
                        content_blocks.append(asset.read_block(locator=page_index))
                    ocr_text = self._run_tesseract_bytes(
                        data=rendered_page,
                        suffix=f"-page-{page_index}.png",
                    )
                    if ocr_text is not None:
                        content_blocks.insert(
                            0,
                            {
                                "type": "text",
                                "locator": page_index,
                                "text": ocr_text.text,
                                "text_source": "ocr",
                            },
                        )
                        search_entries.append(("ocr_text", ocr_text.text))

            if asset_refs and not text and ocr_text is None:
                search_entries.append(
                    (
                        "vision_summary",
                        f"Page {page_index} contains {len(asset_refs)} embedded images",
                    )
                )

            if not content_blocks:
                content_blocks.append(
                    {
                        "type": "document",
                        "locator": page_index,
                        "summary": "No extractable text or image assets on this page.",
                    }
                )

            units.append(
                DocumentUnit(
                    locator=page_index,
                    locator_type="page",
                    content_blocks=content_blocks,
                    search_entries=search_entries,
                    contains_visual=bool(asset_refs),
                    preview_text=text or (ocr_text.text if ocr_text is not None else f"Page {page_index}"),
                )
            )

        return ParsedDocument(
            path=path,
            document_kind="pdf",
            locator_type="page",
            units=units,
            assets=assets,
            contains_visual=bool(assets),
            ingest_metadata=self._default_ingest_metadata(
                ocr_attempted=ocr_attempted,
                ocr_provider="tesseract" if ocr_attempted else None,
            ),
        )

    def _parse_pptx_document(self, *, file_path: Path, path: str) -> ParsedDocument:
        presentation = Presentation(str(file_path))
        units: list[DocumentUnit] = []
        assets: dict[str, DocumentAsset] = {}
        ocr_attempted = False
        for slide_index, slide in enumerate(presentation.slides, start=1):
            text_blocks: list[str] = []
            search_entries: list[tuple[EvidenceType, str]] = []
            content_blocks: list[dict[str, Any]] = []
            contains_visual = False

            for shape_index, shape in enumerate(slide.shapes, start=1):
                if getattr(shape, "has_text_frame", False):
                    text = _collapse_whitespace(getattr(shape, "text", ""))
                    if text:
                        text_blocks.append(text)

                if getattr(shape, "has_table", False):
                    rows = [
                        [_collapse_whitespace(cell.text) for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    region = f"slide:{slide_index}:table:{shape_index}"
                    table_text = _table_rows_to_text(rows)
                    if table_text:
                        search_entries.append(("table_text", table_text))
                    content_blocks.append(
                        {
                            "type": "table",
                            "locator": slide_index,
                            "region": region,
                            "rows": rows,
                            "text": table_text,
                        }
                    )

                if getattr(shape, "has_chart", False):
                    contains_visual = True
                    chart = shape.chart
                    summary = _chart_summary_text(
                        getattr(chart, "chart_title", None),
                        f"Chart on slide {slide_index}",
                    )
                    series_names = [
                        _collapse_whitespace(str(getattr(series, "name", "")))
                        for series in getattr(chart, "series", [])
                        if _collapse_whitespace(str(getattr(series, "name", "")))
                    ]
                    full_summary = "; ".join(
                        part
                        for part in [summary, ", ".join(series_names) if series_names else ""]
                        if part
                    )
                    if full_summary:
                        search_entries.append(("vision_summary", full_summary))
                        content_blocks.append(
                            {
                                "type": "document",
                                "locator": slide_index,
                                "summary": full_summary,
                                "kind": "chart",
                            }
                        )

                image = getattr(shape, "image", None)
                if image is not None:
                    contains_visual = True
                    asset_ref = f"slide:{slide_index}:image:{shape_index}"
                    asset = _asset_from_bytes(
                        asset_ref=asset_ref,
                        data=image.blob,
                        file_name=f"{shape.name}.{image.ext}",
                        summary=f"Slide {slide_index} image {shape_index}",
                        extra={"slide": slide_index, "shape": shape.name, "source": "pptx"},
                    )
                    assets[asset_ref] = asset
                    content_blocks.append(asset.read_block(locator=slide_index))
                    ocr_attempted = True
                    ocr_text = self._ocr_asset_bytes(
                        data=asset.data,
                        asset_ref=asset_ref,
                    )
                    if ocr_text is not None:
                        search_entries.append(("ocr_text", ocr_text.text))
                        content_blocks.append(
                            {
                                "type": "text",
                                "locator": slide_index,
                                "text": ocr_text.text,
                                "text_source": "ocr",
                                "asset_ref": asset_ref,
                            }
                        )

            slide_text = "\n".join(text_blocks).strip()
            if slide_text:
                content_blocks.insert(
                    0,
                    {
                        "type": "text",
                        "locator": slide_index,
                        "text": slide_text,
                    },
                )
                search_entries.insert(0, ("text", slide_text))
            if not content_blocks:
                content_blocks.append(
                    {
                        "type": "document",
                        "locator": slide_index,
                        "summary": f"Slide {slide_index} has no extractable text.",
                    }
                )
            units.append(
                DocumentUnit(
                    locator=slide_index,
                    locator_type="slide",
                    content_blocks=content_blocks,
                    search_entries=search_entries,
                    contains_visual=contains_visual,
                    preview_text=slide_text or f"Slide {slide_index}",
                )
            )

        return ParsedDocument(
            path=path,
            document_kind="pptx",
            locator_type="slide",
            units=units,
            assets=assets,
            contains_visual=any(unit.contains_visual for unit in units),
            ingest_metadata=self._default_ingest_metadata(
                ocr_attempted=ocr_attempted,
                ocr_provider="tesseract" if ocr_attempted else None,
            ),
        )

    def _parse_docx_document(self, *, file_path: Path, path: str) -> ParsedDocument:
        document = WordDocument(str(file_path))
        units: list[DocumentUnit] = []
        assets, asset_refs = _extract_docx_assets(document)
        ocr_attempted = bool(asset_refs)

        for index, block in enumerate(_iter_docx_blocks(document), start=1):
            if isinstance(block, Paragraph):
                text = _collapse_whitespace(block.text)
                content_blocks = [
                    {
                        "type": "text",
                        "locator": index,
                        "text": text,
                    }
                ] if text else [
                    {
                        "type": "document",
                        "locator": index,
                        "summary": f"Region {index} has no visible text.",
                    }
                ]
                units.append(
                    DocumentUnit(
                        locator=index,
                        locator_type="region",
                        content_blocks=content_blocks,
                        search_entries=[("text", text)] if text else [],
                        contains_visual=False,
                        preview_text=text,
                    )
                )
                continue

            rows = [
                [_collapse_whitespace(cell.text) for cell in row.cells]
                for row in block.rows
            ]
            table_text = _table_rows_to_text(rows)
            units.append(
                DocumentUnit(
                    locator=index,
                    locator_type="region",
                    content_blocks=[
                        {
                            "type": "table",
                            "locator": index,
                            "region": f"table:{index}",
                            "rows": rows,
                            "text": table_text,
                        }
                    ],
                    search_entries=[("table_text", table_text)] if table_text else [],
                    contains_visual=False,
                    preview_text=table_text,
                )
            )

        if asset_refs:
            media_locator = len(units) + 1
            media_blocks = [assets[asset_ref].read_block(locator=media_locator) for asset_ref in asset_refs]
            media_entries: list[tuple[EvidenceType, str]] = []
            for asset_ref in asset_refs:
                ocr_text = self._ocr_asset_bytes(
                    data=assets[asset_ref].data,
                    asset_ref=asset_ref,
                )
                if ocr_text is None:
                    continue
                media_entries.append(("ocr_text", ocr_text.text))
                media_blocks.append(
                    {
                        "type": "text",
                        "locator": media_locator,
                        "text": ocr_text.text,
                        "text_source": "ocr",
                        "asset_ref": asset_ref,
                    }
                )
            units.append(
                DocumentUnit(
                    locator=media_locator,
                    locator_type="region",
                    content_blocks=media_blocks,
                    search_entries=[
                        (
                            "vision_summary",
                            f"Document contains {len(asset_refs)} embedded images",
                        ),
                        *media_entries,
                    ],
                    contains_visual=True,
                    preview_text=media_entries[0][1] if media_entries else f"{len(asset_refs)} embedded images",
                )
            )

        if not units:
            units.append(
                DocumentUnit(
                    locator=1,
                    locator_type="region",
                    content_blocks=[
                        {
                            "type": "document",
                            "locator": 1,
                            "summary": "Document has no extractable content.",
                        }
                    ],
                    search_entries=[],
                    contains_visual=bool(asset_refs),
                    preview_text="",
                )
            )

        return ParsedDocument(
            path=path,
            document_kind="docx",
            locator_type="region",
            units=units,
            assets=assets,
            contains_visual=bool(asset_refs),
            ingest_metadata=self._default_ingest_metadata(
                ocr_attempted=ocr_attempted,
                ocr_provider="tesseract" if ocr_attempted else None,
            ),
        )

    def _parse_xlsx_document(self, *, file_path: Path, path: str) -> ParsedDocument:
        workbook = load_workbook(str(file_path), data_only=True)
        try:
            units: list[DocumentUnit] = []
            assets: dict[str, DocumentAsset] = {}
            ocr_attempted = False
            for sheet in workbook.worksheets:
                rows, region = _extract_sheet_rows(sheet)
                table_text = _table_rows_to_text(rows)
                search_entries: list[tuple[EvidenceType, str]] = []
                content_blocks: list[dict[str, Any]] = []
                contains_visual = False

                if rows:
                    content_blocks.append(
                        {
                            "type": "table",
                            "locator": sheet.title,
                            "region": region,
                            "sheet_name": sheet.title,
                            "rows": rows,
                            "text": table_text,
                        }
                    )
                    search_entries.append(("table_text", table_text))

                for chart_index, chart in enumerate(getattr(sheet, "_charts", []), start=1):
                    contains_visual = True
                    summary = _chart_summary_text(
                        getattr(chart, "title", None),
                        f"Chart {chart_index} on sheet {sheet.title}",
                    )
                    search_entries.append(("vision_summary", summary))
                    content_blocks.append(
                        {
                            "type": "document",
                            "locator": sheet.title,
                            "summary": summary,
                            "kind": "chart",
                        }
                    )

                for image_index, image in enumerate(getattr(sheet, "_images", []), start=1):
                    contains_visual = True
                    asset_ref = f"sheet:{sheet.title}:image:{image_index}"
                    assets[asset_ref] = _asset_from_bytes(
                        asset_ref=asset_ref,
                        data=image._data(),
                        file_name=f"{sheet.title}-image-{image_index}.{image.format}",
                        summary=f"Image {image_index} on sheet {sheet.title}",
                        extra={"sheet": sheet.title, "source": "xlsx"},
                    )
                    content_blocks.append(
                        assets[asset_ref].read_block(locator=sheet.title)
                    )
                    ocr_attempted = True
                    ocr_text = self._ocr_asset_bytes(
                        data=assets[asset_ref].data,
                        asset_ref=asset_ref,
                    )
                    if ocr_text is not None:
                        search_entries.append(("ocr_text", ocr_text.text))
                        content_blocks.append(
                            {
                                "type": "text",
                                "locator": sheet.title,
                                "text": ocr_text.text,
                                "text_source": "ocr",
                                "asset_ref": asset_ref,
                            }
                        )

                if not content_blocks:
                    content_blocks.append(
                        {
                            "type": "document",
                            "locator": sheet.title,
                            "summary": f"Sheet {sheet.title} has no extractable content.",
                        }
                    )

                units.append(
                    DocumentUnit(
                        locator=sheet.title,
                        locator_type="sheet",
                        content_blocks=content_blocks,
                        search_entries=search_entries,
                        contains_visual=contains_visual,
                        preview_text=table_text or sheet.title,
                    )
                )
        finally:
            workbook.close()

        return ParsedDocument(
            path=path,
            document_kind="xlsx",
            locator_type="sheet",
            units=units,
            assets=assets,
            contains_visual=any(unit.contains_visual for unit in units),
            ingest_metadata=self._default_ingest_metadata(
                ocr_attempted=ocr_attempted,
                ocr_provider="tesseract" if ocr_attempted else None,
            ),
        )

    def _parse_xls_document(self, *, file_path: Path, path: str) -> ParsedDocument:
        workbook = open_xls_workbook(str(file_path))
        units: list[DocumentUnit] = []
        for sheet in workbook.sheets():
            rows, region = _extract_xls_rows(sheet)
            table_text = _table_rows_to_text(rows)
            content_blocks = [
                {
                    "type": "table",
                    "locator": sheet.name,
                    "region": region,
                    "sheet_name": sheet.name,
                    "rows": rows,
                    "text": table_text,
                }
            ] if rows else [
                {
                    "type": "document",
                    "locator": sheet.name,
                    "summary": f"Sheet {sheet.name} has no extractable content.",
                }
            ]
            units.append(
                DocumentUnit(
                    locator=sheet.name,
                    locator_type="sheet",
                    content_blocks=content_blocks,
                    search_entries=[("table_text", table_text)] if table_text else [],
                    contains_visual=False,
                    preview_text=table_text or sheet.name,
                )
            )
        return ParsedDocument(
            path=path,
            document_kind="xls",
            locator_type="sheet",
            units=units,
            assets={},
            contains_visual=False,
            ingest_metadata=self._default_ingest_metadata(ocr_attempted=False),
        )

    def _parse_legacy_document(self, *, file_path: Path, path: str) -> ParsedDocument:
        kind = LEGACY_OFFICE_EXTENSIONS[file_path.suffix.lower()]
        summary = (
            "Legacy Office binary format. The 8084 demo keeps the document "
            "contract available, but only file-name level search is indexed for "
            "this format because no OLE parser is bundled in the demo image."
        )
        return ParsedDocument(
            path=path,
            document_kind=kind,
            locator_type="region",
            units=[
                DocumentUnit(
                    locator=1,
                    locator_type="region",
                    content_blocks=[
                        {
                            "type": "document",
                            "locator": 1,
                            "summary": summary,
                        }
                    ],
                    search_entries=[("vision_summary", f"{file_path.stem} {summary}")],
                    contains_visual=False,
                    preview_text=summary,
                )
            ],
            assets={},
            contains_visual=False,
            ingest_metadata=self._default_ingest_metadata(ocr_attempted=False),
        )
