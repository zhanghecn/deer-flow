from __future__ import annotations

import sys
import tempfile
import unittest
import json
from pathlib import Path

from docx import Document as WordDocument
from docx.shared import Inches as DocInches
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.drawing.image import Image as XlsxImage
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches as PptxInches


SERVICE_ROOT = Path(__file__).resolve().parents[1]
if str(SERVICE_ROOT) not in sys.path:
    sys.path.insert(0, str(SERVICE_ROOT))

from app.service import FileMcpService


def build_pdf_bytes(pages: list[str]) -> bytes:
    """Generate a minimal text PDF that pypdf can extract in tests."""

    header = b"%PDF-1.4\n"
    objects: list[bytes] = [
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n",
    ]
    page_refs: list[str] = []
    page_objects: list[bytes] = []
    next_object_number = 3
    font_object_number = 2 * len(pages) + 3

    for text in pages:
        page_number = next_object_number
        content_number = next_object_number + 1
        page_refs.append(f"{page_number} 0 R")
        safe_text = (
            text.replace("\\", "\\\\")
            .replace("(", "\\(")
            .replace(")", "\\)")
        )
        stream = (
            f"BT\n/F1 18 Tf\n72 720 Td\n({safe_text}) Tj\nET\n".encode("latin-1")
        )
        page_objects.append(
            (
                f"{page_number} 0 obj\n"
                f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
                f"/Resources << /Font << /F1 {font_object_number} 0 R >> >> "
                f"/Contents {content_number} 0 R >>\nendobj\n"
            ).encode("latin-1")
        )
        page_objects.append(
            (
                f"{content_number} 0 obj\n<< /Length {len(stream)} >>\nstream\n"
            ).encode("latin-1")
            + stream
            + b"endstream\nendobj\n"
        )
        next_object_number += 2

    objects.append(
        (
            f"2 0 obj\n<< /Type /Pages /Kids [{' '.join(page_refs)}] "
            f"/Count {len(page_refs)} >>\nendobj\n"
        ).encode("latin-1")
    )
    objects.extend(page_objects)
    objects.append(
        (
            f"{font_object_number} 0 obj\n"
            "<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica "
            "/Encoding /WinAnsiEncoding >>\nendobj\n"
        ).encode("latin-1")
    )

    offsets = [0]
    body = b""
    for obj in objects:
        offsets.append(len(header) + len(body))
        body += obj

    startxref = len(header) + len(body)
    xref_parts = [
        f"xref\n0 {len(offsets)}\n0000000000 65535 f \n".encode("latin-1")
    ]
    for offset in offsets[1:]:
        xref_parts.append(f"{offset:010d} 00000 n \n".encode("latin-1"))

    trailer = (
        f"trailer\n<< /Size {len(offsets)} /Root 1 0 R >>\n"
        f"startxref\n{startxref}\n%%EOF\n"
    ).encode("latin-1")
    return header + body + b"".join(xref_parts) + trailer


class FileMcpServiceTest(unittest.TestCase):
    def setUp(self) -> None:
        self.temp_dir = tempfile.TemporaryDirectory()
        self.root = Path(self.temp_dir.name)
        cases_dir = self.root / "案例大全"
        cases_dir.mkdir(parents=True, exist_ok=True)
        (cases_dir / "a.md").write_text(
            "灾祸 会来\n好运 也会来\n官非 风险\n",
            encoding="utf-8",
        )
        (cases_dir / "b.md").write_text(
            "血光 之灾\n顺遂 发展\n",
            encoding="utf-8",
        )
        self.service = FileMcpService(
            self.root,
            ocr_languages=("eng",),
            public_base_url="http://127.0.0.1:8084",
        )

    def tearDown(self) -> None:
        self.temp_dir.cleanup()

    def _write_image(self, relative_path: str, *, text: str | None = None) -> Path:
        target = self.root / relative_path
        target.parent.mkdir(parents=True, exist_ok=True)
        image = Image.new("RGB", (1200, 320), "white")
        if text:
            draw = ImageDraw.Draw(image)
            font_path = Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
            if font_path.exists():
                font = ImageFont.truetype(str(font_path), 96)
            else:
                font = ImageFont.load_default()
            draw.text((60, 90), text, fill="black", font=font)
        image.save(target)
        return target

    def _write_scanned_pdf(self, relative_path: str, *, text: str) -> Path:
        image_path = self._write_image("assets/scanned-page.png", text=text)
        target = self.root / relative_path
        target.parent.mkdir(parents=True, exist_ok=True)
        with Image.open(image_path) as image:
            image.convert("RGB").save(target, "PDF", resolution=150)
        return target

    def _write_pdf(self, relative_path: str, pages: list[str]) -> Path:
        target = self.root / relative_path
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_bytes(build_pdf_bytes(pages))
        return target

    def _write_text(self, relative_path: str, content: str) -> Path:
        target = self.root / relative_path
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(content, encoding="utf-8")
        return target

    def _write_pptx(self, relative_path: str) -> Path:
        image_path = self._write_image("assets/slide-proof.png")
        target = self.root / relative_path
        target.parent.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Quarterly Revenue"
        slide.shapes.add_picture(
            str(image_path),
            PptxInches(0.8),
            PptxInches(1.2),
            width=PptxInches(1.8),
        )
        chart_data = CategoryChartData()
        chart_data.categories = ["Q1", "Q2"]
        chart_data.add_series("Revenue", (120, 180))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            PptxInches(3.0),
            PptxInches(1.2),
            PptxInches(4.5),
            PptxInches(3.2),
            chart_data,
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "Revenue Trend"
        presentation.save(target)
        return target

    def _write_pptx_with_ocr_image(self, relative_path: str, *, text: str) -> Path:
        image_path = self._write_image("assets/slide-ocr.png", text=text)
        target = self.root / relative_path
        target.parent.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(
            str(image_path),
            PptxInches(0.8),
            PptxInches(0.8),
            width=PptxInches(8.0),
        )
        presentation.save(target)
        return target

    def _write_docx(self, relative_path: str) -> Path:
        image_path = self._write_image("assets/doc-proof.png")
        target = self.root / relative_path
        target.parent.mkdir(parents=True, exist_ok=True)
        document = WordDocument()
        document.add_paragraph("Customer premium note")
        document.add_picture(str(image_path), width=DocInches(1))
        document.save(target)
        return target

    def _write_xlsx(self, relative_path: str) -> Path:
        image_path = self._write_image("assets/sheet-proof.png")
        target = self.root / relative_path
        target.parent.mkdir(parents=True, exist_ok=True)
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Revenue"
        sheet.append(["Quarter", "Revenue"])
        sheet.append(["Q1", 120])
        sheet.append(["Q2", 180])
        sheet.add_image(XlsxImage(str(image_path)), "D2")
        chart = BarChart()
        chart.title = "Revenue Trend"
        data = Reference(sheet, min_col=2, min_row=1, max_row=3)
        categories = Reference(sheet, min_col=1, min_row=2, max_row=3)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(categories)
        sheet.add_chart(chart, "F2")
        workbook.save(target)
        workbook.close()
        return target

    def test_list_reports_binary_documents_without_hiding_sibling_text_files(self) -> None:
        (self.root / "合同.pdf").write_bytes(b"%PDF-1.4 fake")
        (self.root / "合同.md").write_text(
            "# 合同\n\n赔付条款\n",
            encoding="utf-8",
        )

        payload = self.service.list_files_payload(limit=20)
        items_by_path = {item["path"]: item for item in payload["items"]}

        self.assertIn("合同.pdf", items_by_path)
        self.assertIn("合同.md", items_by_path)
        self.assertEqual(items_by_path["合同.pdf"]["content_kind"], "binary_document")
        self.assertFalse(items_by_path["合同.pdf"]["text_readable"])
        self.assertEqual(items_by_path["合同.md"]["content_kind"], "text")
        self.assertTrue(items_by_path["合同.md"]["text_readable"])

    def test_workspace_virtual_root_alias_resolves_to_uploaded_root(self) -> None:
        payload = self.service.ls_payload(path="/mnt/user-data/workspace", limit=20)
        listed_paths = {item["path"] for item in payload["items"]}

        self.assertIn("案例大全", listed_paths)
        self.assertEqual(payload["total"], 1)

    def test_document_list_returns_final_files_with_document_metadata(self) -> None:
        self._write_pdf(
            "nested/contracts/policy-alpha.pdf",
            ["Deductible is 500 USD"],
        )
        self._write_image("nested/evidence/site-board.png", text="POLICY LIMIT")
        self._write_pdf("root-policy.pdf", ["Root policy"])

        root_payload = self.service.document_list_payload()
        root_lines = root_payload["content"].splitlines()
        self.assertIn("- nested/contracts/policy-alpha.pdf [pdf]", root_payload["content"])
        self.assertIn("- nested/evidence/site-board.png [image visual]", root_payload["content"])
        self.assertIn("- root-policy.pdf [pdf]", root_payload["content"])
        self.assertNotIn("- nested/", root_lines)
        self.assertNotIn("- nested/contracts/", root_lines)
        self.assertEqual(root_payload["total"], len(root_payload["items"]))
        self.assertEqual(root_payload["returned"], root_payload["total"])
        self.assertFalse(root_payload["has_more"])
        items_by_path = {item["path"]: item for item in root_payload["items"]}
        self.assertNotIn("nested", items_by_path)
        self.assertNotIn("nested/contracts", items_by_path)
        self.assertEqual(items_by_path["root-policy.pdf"]["document_kind"], "pdf")

        nested_payload = self.service.document_list_payload(path="nested/contracts")
        self.assertIn("- nested/contracts/policy-alpha.pdf [pdf]", nested_payload["content"])
        self.assertEqual(nested_payload["total"], 1)

    def test_document_list_returns_all_tree_entries_without_pagination(self) -> None:
        for index in range(351):
            target = self.root / "bulk" / f"doc-{index:03d}.md"
            target.parent.mkdir(parents=True, exist_ok=True)
            target.write_text(f"# Doc {index}\n", encoding="utf-8")

        payload = self.service.document_list_payload(path="bulk")

        self.assertEqual(payload["total"], 351)
        self.assertEqual(payload["returned"], 351)
        self.assertFalse(payload["has_more"])
        self.assertIn("- bulk/doc-000.md [text]", payload["content"])
        self.assertIn("- bulk/doc-350.md [text]", payload["content"])
        self.assertNotIn("next_offset", payload)

    def test_read_rejects_binary_documents_instead_of_converting_them(self) -> None:
        (self.root / "合同.pdf").write_bytes(b"%PDF-1.4 fake")

        with self.assertRaisesRegex(
            ValueError,
            "fs_read only supports text files.*document_\\* MCP tools",
        ):
            self.service.read_file_payload(file_path="合同.pdf", offset=0, limit=10)

    def test_preview_explains_binary_documents_without_fake_text_content(self) -> None:
        (self.root / "合同.pdf").write_bytes(b"%PDF-1.4 fake")

        payload = self.service.preview_file_payload(path="合同.pdf")

        self.assertEqual(payload["content_kind"], "binary_document")
        self.assertFalse(payload["text_readable"])
        self.assertIn("document_list(path?)", payload["content"])

    def test_grep_skips_binary_documents_in_directory_scope(self) -> None:
        (self.root / "合同.pdf").write_bytes(b"%PDF-1.4 fake")

        payload = self.service.grep_payload(pattern="赔付", path="", glob="*.pdf", limit=20)

        self.assertEqual(payload["total"], 0)
        self.assertEqual(payload["skipped_binary_files"], 1)

    def test_grep_rejects_binary_document_file_scope(self) -> None:
        (self.root / "合同.pdf").write_bytes(b"%PDF-1.4 fake")

        with self.assertRaisesRegex(
            ValueError,
            "fs_grep only supports text files.*document_\\* MCP tools",
        ):
            self.service.grep_payload(pattern="赔付", path="合同.pdf", limit=20)

    def test_delete_file_only_removes_requested_path(self) -> None:
        original = self.root / "合同.pdf"
        companion = self.root / "合同.md"
        original.write_bytes(b"%PDF-1.4 fake")
        companion.write_text("# 合同\n", encoding="utf-8")

        self.service.delete_file("合同.pdf")

        self.assertFalse(original.exists())
        self.assertTrue(companion.exists())

    def test_grep_accepts_count_output_mode(self) -> None:
        payload = self.service.grep_payload(
            pattern="灾祸|血光|官非",
            path="/mnt/user-data/uploads/案例大全",
            glob="*.md",
            output_mode="count",
            limit=20,
        )

        self.assertEqual(payload["output_mode"], "count")
        self.assertEqual(payload["requested_output_mode"], "count")
        self.assertEqual(payload["total_matches"], 3)
        self.assertEqual(
            payload["items"],
            [
                {"path": "案例大全/a.md", "match_count": 2},
                {"path": "案例大全/b.md", "match_count": 1},
            ],
        )

    def test_grep_treats_regex_alternation_as_matches(self) -> None:
        payload = self.service.grep_payload(
            pattern="灾祸|血光|官非",
            path="/mnt/user-data/uploads/案例大全",
            glob="*.md",
            output_mode="content",
            limit=20,
        )

        self.assertEqual(payload["output_mode"], "content")
        self.assertEqual(payload["total"], 3)
        self.assertEqual(
            [item["path"] for item in payload["items"]],
            ["案例大全/a.md", "案例大全/a.md", "案例大全/b.md"],
        )

    def test_grep_normalizes_file_alias_mode(self) -> None:
        payload = self.service.grep_payload(
            pattern="好运|顺遂",
            path="/mnt/user-data/uploads/案例大全",
            glob="*.md",
            output_mode="files",
            limit=20,
        )

        self.assertEqual(payload["output_mode"], "files_with_matches")
        self.assertEqual(payload["requested_output_mode"], "files")
        self.assertEqual(payload["items"], ["案例大全/a.md", "案例大全/b.md"])

    def test_document_search_greps_pdf_text(self) -> None:
        self._write_pdf(
            "nested/contracts/policy-alpha.pdf",
            ["Deductible is 500 USD", "Coinsurance is 20 percent"],
        )

        payload = self.service.document_search_payload(
            pattern="Deductible",
            path="nested/contracts/policy-alpha.pdf",
            output_mode="content",
            head_limit=5,
        )

        self.assertEqual(payload["total"], 1)
        self.assertEqual(payload["mode"], "content")
        self.assertEqual(payload["pattern"], "Deductible")
        match = payload["matches"][0]
        self.assertEqual(match["document_kind"], "pdf")
        self.assertEqual(match["locator_type"], "page")
        self.assertEqual(match["locator"], 1)
        self.assertEqual(match["evidence_type"], "text")
        self.assertIn("Deductible", payload["content"])
        self.assertEqual(match["next_action_hint"], "read_more")
        self.assertEqual(match["read_args"]["locator"], 1)
        self.assertEqual(match["read_args"]["path"], "nested/contracts/policy-alpha.pdf")

    def test_document_search_does_not_decompose_natural_language_questions(self) -> None:
        self._write_pdf(
            "nested/contracts/policy-alpha.pdf",
            ["Deductible is 500 USD", "Coinsurance is 20 percent"],
        )

        payload = self.service.document_search_payload(
            pattern="What is the deductible in policy alpha?",
            path="nested/contracts/policy-alpha.pdf",
            output_mode="content",
            head_limit=5,
        )

        self.assertEqual(payload["mode"], "content")
        self.assertEqual(payload["total"], 0)

    def test_document_search_defaults_to_content_with_line_numbers(self) -> None:
        self._write_pdf(
            "nested/contracts/policy-alpha.pdf",
            ["Deductible is 500 USD", "Coinsurance is 20 percent"],
        )

        payload = self.service.document_search_payload(
            pattern="Deductible",
            path="nested",
        )

        self.assertEqual(payload["mode"], "content")
        self.assertIn("nested/contracts/policy-alpha.pdf:page:1:1", payload["content"])
        self.assertIn("Deductible is 500 USD", payload["content"])
        self.assertEqual(payload["numFiles"], 1)

    def test_document_search_keeps_files_with_matches_mode(self) -> None:
        self._write_pdf(
            "nested/contracts/policy-alpha.pdf",
            ["Deductible is 500 USD", "Coinsurance is 20 percent"],
        )

        payload = self.service.document_search_payload(
            pattern="Deductible",
            path="nested",
            output_mode="files_with_matches",
        )

        self.assertEqual(payload["mode"], "files_with_matches")
        self.assertEqual(payload["filenames"], ["nested/contracts/policy-alpha.pdf"])
        self.assertEqual(payload["numFiles"], 1)

    def test_document_search_supports_regex_glob_context_and_offset(self) -> None:
        self._write_pdf(
            "nested/contracts/policy-alpha.pdf",
            ["Header", "Deductible is 500 USD", "Coinsurance is 20 percent"],
        )
        self._write_pdf(
            "nested/contracts/policy-beta.pdf",
            ["Header", "Deductible is 250 USD", "Coinsurance is 10 percent"],
        )

        payload = self.service.document_search_payload(
            pattern="Deductible|Coinsurance",
            path="nested",
            glob="**/policy-*.pdf",
            output_mode="content",
            context=1,
            head_limit=1,
            offset=1,
        )

        self.assertEqual(payload["mode"], "content")
        self.assertEqual(payload["appliedLimit"], 1)
        self.assertEqual(payload["appliedOffset"], 1)
        self.assertIn("nested/contracts/policy-alpha.pdf:page:3:1", payload["content"])
        self.assertIn("Coinsurance", payload["content"])
        self.assertNotIn("context_lines", payload["matches"][0])

    def test_document_search_content_output_is_context_safe(self) -> None:
        self._write_text(
            "cases/large.md",
            "\n".join(
                [
                    f"prefix {index} 甲辰 " + ("detail " * 80)
                    for index in range(80)
                ]
            ),
        )

        payload = self.service.document_search_payload(
            pattern="甲辰",
            path="cases/large.md",
            output_mode="content",
            context=15,
            head_limit=50,
        )

        self.assertEqual(payload["mode"], "content")
        self.assertNotIn("context_truncated", payload)
        self.assertEqual(payload["applied_context_before"], 15)
        self.assertEqual(payload["applied_context_after"], 15)
        self.assertLessEqual(len(payload["matches"]), 12)
        self.assertGreater(len(payload["matches"]), 0)
        self.assertGreater(payload["numLines"], 0)
        self.assertIn("read_args", payload["matches"][0])
        self.assertTrue(payload["has_more"])
        self.assertNotIn("context_lines", payload["matches"][0])
        self.assertLessEqual(
            len(json.dumps(payload, ensure_ascii=False).encode("utf-8")),
            8_000,
        )

    def test_document_search_content_preserves_numeric_line_order(self) -> None:
        self._write_text(
            "cases/order.md",
            "\n".join(
                "needle" if index in {3, 10, 100} else f"line {index}"
                for index in range(1, 105)
            ),
        )

        payload = self.service.document_search_payload(
            pattern="needle",
            path="cases/order.md",
            output_mode="content",
            head_limit=3,
        )

        lines = payload["content"].splitlines()
        self.assertIn("cases/order.md:5:needle", lines[0])
        self.assertIn("cases/order.md:12:needle", lines[1])
        self.assertIn("cases/order.md:102:needle", lines[2])

    def test_document_search_count_mode_returns_file_counts(self) -> None:
        self._write_pdf(
            "nested/contracts/policy-alpha.pdf",
            ["Deductible", "Deductible", "Coinsurance"],
        )

        payload = self.service.document_search_payload(
            pattern="Deductible",
            path="nested",
            output_mode="count",
        )

        self.assertEqual(payload["mode"], "count")
        self.assertEqual(payload["numMatches"], 2)
        self.assertEqual(
            payload["counts"],
            [{"path": "nested/contracts/policy-alpha.pdf", "match_count": 2}],
        )

    def test_document_read_paginates_cached_markdown_lines(self) -> None:
        self._write_pdf(
            "nested/contracts/policy-alpha.pdf",
            ["Deductible is 500 USD", "Coinsurance is 20 percent"],
        )

        payload = self.service.document_read_payload(
            path="nested/contracts/policy-alpha.pdf",
            offset=0,
            limit=1,
        )

        self.assertEqual(payload["document_kind"], "pdf")
        self.assertEqual(payload["locator_type"], "page")
        self.assertEqual(payload["total_units"], 2)
        self.assertEqual(payload["read_surface"], "canonical_markdown")
        self.assertEqual(payload["total_lines"], 13)
        self.assertEqual(payload["returned_lines"], 1)
        self.assertTrue(payload["has_more"])
        self.assertEqual(payload["next_offset"], 1)
        self.assertFalse(payload["contains_visual"])
        self.assertEqual(
            payload["source_url"],
            "http://127.0.0.1:8084/api/files/source?path=nested%2Fcontracts%2Fpolicy-alpha.pdf#page=1",
        )
        self.assertEqual(payload["content"], "# nested/contracts/policy-alpha.pdf\n")
        full_payload = self.service.document_read_payload(
            path="nested/contracts/policy-alpha.pdf",
        )
        self.assertIn("## Page 1", full_payload["content"])
        self.assertIn("Deductible is 500 USD", full_payload["content"])

        cache_dir = self.service.cache_root / "nested" / "contracts" / "policy-alpha.pdf"
        self.assertTrue((cache_dir / "manifest.json").exists())
        self.assertTrue((cache_dir / "canonical.md").exists())
        self.assertTrue((cache_dir / "parse_result.json").exists())
        self.assertIn(
            "Deductible is 500 USD",
            (cache_dir / "canonical.md").read_text(encoding="utf-8"),
        )
        parse_result = json.loads((cache_dir / "parse_result.json").read_text(encoding="utf-8"))
        self.assertEqual(parse_result["parse_mode"], "local")
        self.assertEqual(parse_result["layout_visualizations"], {})
        self.assertIn("Deductible is 500 USD", parse_result["markdown"][0])
        self.assertEqual(payload["local_parse"]["parse_result_path"], "parse_result.json")

    def test_document_read_caps_large_text_windows(self) -> None:
        self._write_text(
            "cases/long.md",
            "\n".join(
                [
                    f"case {index} 甲辰 " + ("detail " * 2000)
                    for index in range(80)
                ]
            ),
        )

        payload = self.service.document_read_payload(
            path="cases/long.md",
            offset=0,
            limit=50,
        )

        self.assertEqual(payload["requested_limit"], 50)
        self.assertEqual(payload["limit"], 50)
        self.assertEqual(payload["read_surface"], "canonical_markdown")
        self.assertTrue(payload["content_truncated"])
        self.assertTrue(payload["has_more"])
        self.assertLess(payload["returned_lines"], 50)
        self.assertIn("[truncated]", payload["content"])
        self.assertLessEqual(
            len(json.dumps(payload, ensure_ascii=False).encode("utf-8")),
            70_000,
        )

    def test_document_read_uses_search_locator_without_cursor_guessing(self) -> None:
        self._write_text(
            "cases/sections.md",
            "\n".join(
                [
                    "案例一 header",
                    "普通内容",
                    "案例二 header",
                    "八字 甲辰 甲戌 戊子 甲寅",
                    "真实事件 detail",
                    "命理分析 detail",
                    "结尾",
                ]
            ),
        )

        search_payload = self.service.document_search_payload(
            pattern="甲辰 甲戌",
            path="cases/sections.md",
            output_mode="content",
            head_limit=1,
        )
        read_args = search_payload["matches"][0]["read_args"]
        self.assertEqual(
            search_payload["source_links"][0]["source_url"],
            "http://127.0.0.1:8084/api/files/source?path=cases%2Fsections.md",
        )
        read_payload = self.service.document_read_payload(
            **read_args,
            limit=12,
        )

        self.assertEqual(read_payload["locator"], "4")
        self.assertEqual(
            read_payload["source_url"],
            "http://127.0.0.1:8084/api/files/source?path=cases%2Fsections.md&locator=4",
        )
        self.assertEqual(read_payload["matched_unit_offset"], 3)
        self.assertGreaterEqual(read_payload["matched_offset"], 3)
        rendered = read_payload["content"]
        self.assertIn("案例二 header", rendered)
        self.assertIn("八字 甲辰 甲戌 戊子 甲寅", rendered)
        self.assertIn("命理分析 detail", rendered)

    def test_document_search_ocr_hits_scanned_pdf(self) -> None:
        if not self.service.document_tools.ocr_available:
            self.skipTest("tesseract is not available in this test environment")

        self._write_scanned_pdf("nested/scans/claim-scan.pdf", text="CLAIM 42")

        search_payload = self.service.document_search_payload(
            pattern="CLAIM",
            path="nested/scans/claim-scan.pdf",
            output_mode="content",
            head_limit=5,
        )
        read_payload = self.service.document_read_payload(
            path="nested/scans/claim-scan.pdf",
            offset=0,
            limit=1,
        )

        self.assertEqual(search_payload["matches"][0]["evidence_type"], "ocr_text")
        self.assertTrue(search_payload["matches"][0]["contains_visual"])
        self.assertTrue(read_payload["contains_visual"])
        self.assertIn("[text source=ocr]", read_payload["content"])
        self.assertTrue(read_payload["assets"])

        manifest = (
            self.service.cache_root / "nested" / "scans" / "claim-scan.pdf" / "manifest.json"
        ).read_text(encoding="utf-8")
        self.assertIn('"ocr_status": "complete"', manifest)

    def test_document_read_returns_pptx_text_and_visual_blocks(self) -> None:
        self._write_pptx("nested/slides/review-deck.pptx")

        payload = self.service.document_read_payload(
            path="nested/slides/review-deck.pptx",
        )

        self.assertEqual(payload["document_kind"], "pptx")
        self.assertTrue(payload["contains_visual"])
        self.assertEqual(payload["read_surface"], "canonical_markdown")
        self.assertIn("## Slide 1", payload["content"])
        self.assertIn("![", payload["content"])
        self.assertTrue(payload["image_paths"])
        self.assertEqual(
            payload["image_read_args"][0],
            {"path": "nested/slides/review-deck.pptx", "locator": payload["image_paths"][0]},
        )

        cache_dir = self.service.cache_root / "nested" / "slides" / "review-deck.pptx"
        cached_images = sorted((cache_dir / "images").iterdir())
        self.assertTrue(cached_images)
        self.assertTrue(any(item.name.startswith("slide_1_image_2") for item in cached_images))
        parse_result = json.loads((cache_dir / "parse_result.json").read_text(encoding="utf-8"))
        self.assertTrue(parse_result["image_mapping"])
        self.assertTrue(parse_result["image_paths"])
        image_path = parse_result["image_paths"][0]
        self.assertTrue(image_path.startswith("images/"))
        self.assertIn(f"]({image_path})", parse_result["markdown"][0])

    def test_document_search_ocr_hits_pptx_image_slide(self) -> None:
        if not self.service.document_tools.ocr_available:
            self.skipTest("tesseract is not available in this test environment")

        self._write_pptx_with_ocr_image(
            "nested/slides/ocr-deck.pptx",
            text="LOSS RATIO",
        )

        payload = self.service.document_search_payload(
            pattern="LOSS",
            path="nested/slides/ocr-deck.pptx",
            output_mode="content",
            head_limit=5,
        )

        self.assertEqual(payload["matches"][0]["document_kind"], "pptx")
        self.assertEqual(payload["matches"][0]["evidence_type"], "ocr_text")
        self.assertEqual(payload["matches"][0]["next_action_hint"], "read_current")

    def test_prime_document_cache_builds_existing_seed_documents(self) -> None:
        self._write_docx("nested/docs/briefing.docx")
        self._write_xlsx("nested/sheets/revenue-tracker.xlsx")

        self.service.prime_document_cache()

        self.assertTrue((self.service.cache_root / "nested" / "docs" / "briefing.docx" / "manifest.json").exists())
        self.assertTrue((self.service.cache_root / "nested" / "sheets" / "revenue-tracker.xlsx" / "canonical.md").exists())

    def test_document_read_locator_inlines_small_pptx_images(self) -> None:
        self._write_pptx("nested/slides/review-deck.pptx")
        read_payload = self.service.document_read_payload(
            path="nested/slides/review-deck.pptx",
        )
        image_path = read_payload["image_paths"][0]

        payload = self.service.document_read_payload(
            path="nested/slides/review-deck.pptx",
            locator=image_path,
        )

        self.assertEqual(payload["read_surface"], "cached_image")
        self.assertEqual(payload["image_paths"], [image_path])
        self.assertTrue(payload["mcp_images"])
        self.assertIn("content_base64", payload["mcp_images"][0])

    def test_document_read_locator_rejects_missing_image_path(self) -> None:
        self._write_pdf("nested/contracts/policy-alpha.pdf", ["No images here"])

        with self.assertRaises(FileNotFoundError):
            self.service.document_read_payload(
                path="nested/contracts/policy-alpha.pdf",
                locator="images/missing.png",
            )

    def test_document_search_reads_xlsx_table_and_chart_evidence(self) -> None:
        self._write_xlsx("nested/sheets/revenue-tracker.xlsx")

        table_payload = self.service.document_search_payload(
            pattern="Q2",
            path="nested/sheets/revenue-tracker.xlsx",
            output_mode="content",
            head_limit=5,
        )
        chart_payload = self.service.document_search_payload(
            pattern="Revenue Trend",
            path="nested/sheets/revenue-tracker.xlsx",
            output_mode="content",
            head_limit=5,
        )

        self.assertEqual(table_payload["matches"][0]["evidence_type"], "table_text")
        self.assertEqual(chart_payload["matches"][0]["evidence_type"], "vision_summary")
        self.assertEqual(chart_payload["matches"][0]["locator_type"], "sheet")

    def test_document_read_returns_xlsx_table_block_and_visual_flag(self) -> None:
        self._write_xlsx("nested/sheets/revenue-tracker.xlsx")

        payload = self.service.document_read_payload(
            path="nested/sheets/revenue-tracker.xlsx",
        )

        self.assertEqual(payload["document_kind"], "xlsx")
        self.assertTrue(payload["contains_visual"])
        self.assertEqual(payload["locator_type"], "sheet")
        self.assertIn("```table", payload["content"])

    def test_document_read_returns_docx_regions_and_media_unit(self) -> None:
        self._write_docx("nested/docs/briefing.docx")

        payload = self.service.document_read_payload(
            path="nested/docs/briefing.docx",
            offset=0,
        )

        self.assertEqual(payload["document_kind"], "docx")
        self.assertGreaterEqual(payload["total_units"], 2)
        self.assertTrue(payload["assets"])

    def test_document_search_ocr_hits_image_file(self) -> None:
        if not self.service.document_tools.ocr_available:
            self.skipTest("tesseract is not available in this test environment")

        self._write_image("nested/images/policy-board.png", text="POLICY LIMIT")

        payload = self.service.document_search_payload(
            pattern="POLICY",
            path="nested/images/policy-board.png",
            output_mode="content",
            head_limit=5,
        )

        self.assertEqual(payload["matches"][0]["document_kind"], "image")
        self.assertEqual(payload["matches"][0]["evidence_type"], "ocr_text")
        read_payload = self.service.document_read_payload(
            path="nested/images/policy-board.png",
            offset=0,
            limit=1,
        )
        self.assertTrue(read_payload["contains_visual"])
        self.assertIn("[text source=ocr]", read_payload["content"])


if __name__ == "__main__":
    unittest.main()
