#!/usr/bin/env python3

import json
import os
from html import escape
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _slide_dimensions(aspect_ratio: str) -> tuple[object, object]:
    if aspect_ratio == "4:3":
        return Inches(10), Inches(7.5)
    return Inches(13.333), Inches(7.5)


def _add_notes(slide, slide_info: dict) -> None:
    notes: list[str] = []

    if slide_info.get("title"):
        notes.append(f"Title: {slide_info['title']}")

    if slide_info.get("subtitle"):
        notes.append(f"Subtitle: {slide_info['subtitle']}")

    if slide_info.get("key_points"):
        notes.append("Key Points:")
        for point in slide_info["key_points"]:
            notes.append(f"  • {point}")

    if not notes:
        return

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    if text_frame is not None:
        text_frame.text = "\n".join(notes)


def _add_full_slide_image(slide, image_path: str, slide_width, slide_height) -> None:
    with Image.open(image_path) as img:
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")

        img_width, img_height = img.size
        img_aspect = img_width / img_height
        slide_aspect = slide_width / slide_height

        slide_width_emu = int(slide_width)
        slide_height_emu = int(slide_height)

        if img_aspect > slide_aspect:
            new_width_emu = slide_width_emu
            new_height_emu = int(slide_width_emu / img_aspect)
            left = Inches(0)
            top = Inches((slide_height_emu - new_height_emu) / 914400)
        else:
            new_height_emu = slide_height_emu
            new_width_emu = int(slide_height_emu * img_aspect)
            left = Inches((slide_width_emu - new_width_emu) / 914400)
            top = Inches(0)

        img_bytes = BytesIO()
        img.save(img_bytes, format="JPEG", quality=95)
        img_bytes.seek(0)
        slide.shapes.add_picture(
            img_bytes,
            left,
            top,
            Inches(new_width_emu / 914400),
            Inches(new_height_emu / 914400),
        )


def _set_background(slide, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color_hex)


def _add_text_slide(slide, slide_info: dict, slide_width, slide_height) -> None:
    _set_background(slide, "0F172A")

    title = str(slide_info.get("title") or "Untitled Slide")
    subtitle = str(slide_info.get("subtitle") or "")
    key_points = [str(point) for point in slide_info.get("key_points") or []]
    visual_description = str(slide_info.get("visual_description") or "").strip()

    title_box = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(0.6),
        slide_width - Inches(1.6),
        Inches(1.2),
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor.from_string("F8FAFC")

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(1.55),
            slide_width - Inches(1.6),
            Inches(0.7),
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.text = subtitle
        subtitle_paragraph.font.size = Pt(14)
        subtitle_paragraph.font.color.rgb = RGBColor.from_string("CBD5E1")

    body_box = slide.shapes.add_textbox(
        Inches(0.9),
        Inches(2.4),
        slide_width - Inches(1.8),
        slide_height - Inches(3.0),
    )
    body_frame = body_box.text_frame
    body_frame.word_wrap = True
    body_frame.clear()

    if key_points:
        for index, point in enumerate(key_points):
            paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            paragraph.text = point
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = RGBColor.from_string("E2E8F0")
            paragraph.space_after = Pt(10)
    elif visual_description:
        paragraph = body_frame.paragraphs[0]
        paragraph.text = visual_description
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor.from_string("E2E8F0")
    else:
        paragraph = body_frame.paragraphs[0]
        paragraph.text = "Content intentionally omitted."
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor.from_string("94A3B8")

    footer_box = slide.shapes.add_textbox(
        Inches(0.9),
        slide_height - Inches(0.6),
        slide_width - Inches(1.8),
        Inches(0.3),
    )
    footer_frame = footer_box.text_frame
    footer_frame.word_wrap = False
    footer_paragraph = footer_frame.paragraphs[0]
    footer_paragraph.text = "Generated fallback layout"
    footer_paragraph.alignment = PP_ALIGN.RIGHT
    footer_paragraph.font.size = Pt(10)
    footer_paragraph.font.color.rgb = RGBColor.from_string("64748B")


def _render_preview_html(plan: dict, slide_images: list[str], output_file: str) -> Path:
    output_path = Path(output_file)
    preview_path = output_path.with_suffix(".html")
    slides = plan.get("slides", [])
    image_map = {index: Path(image_path).name for index, image_path in enumerate(slide_images)}
    sections: list[str] = []

    for index, slide_info in enumerate(slides):
        title = escape(str(slide_info.get("title") or f"Slide {index + 1}"))
        subtitle = escape(str(slide_info.get("subtitle") or ""))
        key_points = [escape(str(point)) for point in slide_info.get("key_points") or []]
        image_name = image_map.get(index)

        if image_name:
            body = f'<img src="{escape(image_name)}" alt="{title}" />'
        else:
            bullet_html = "".join(f"<li>{point}</li>" for point in key_points)
            description = escape(str(slide_info.get("visual_description") or ""))
            body = (
                f"<div class=\"text-slide\">"
                f"{f'<p class=\"subtitle\">{subtitle}</p>' if subtitle else ''}"
                f"{f'<ul>{bullet_html}</ul>' if bullet_html else ''}"
                f"{f'<p class=\"description\">{description}</p>' if description else ''}"
                f"</div>"
            )

        sections.append(
            "<section class=\"slide\">"
            f"<div class=\"label\">Slide {index + 1}</div>"
            f"<h2>{title}</h2>"
            f"{body}"
            "</section>"
        )

    html = f"""<!doctype html>
<html lang="zh-CN">
  <head>
    <meta charset="utf-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1" />
    <title>{escape(str(plan.get("title") or output_path.stem))}</title>
    <style>
      :root {{
        color-scheme: light;
        --bg: #f8fafc;
        --card: #ffffff;
        --ink: #0f172a;
        --muted: #64748b;
        --line: #cbd5e1;
        --accent: #2563eb;
      }}
      * {{ box-sizing: border-box; }}
      body {{
        margin: 0;
        padding: 32px 24px 48px;
        background: radial-gradient(circle at top, #dbeafe, var(--bg) 35%);
        color: var(--ink);
        font-family: "Noto Sans SC", "PingFang SC", "Microsoft YaHei", sans-serif;
      }}
      main {{
        max-width: 1120px;
        margin: 0 auto;
      }}
      h1 {{
        margin: 0 0 12px;
        font-size: 36px;
      }}
      p.meta {{
        margin: 0 0 32px;
        color: var(--muted);
      }}
      section.slide {{
        margin-bottom: 28px;
        padding: 24px;
        border: 1px solid var(--line);
        border-radius: 24px;
        background: rgba(255, 255, 255, 0.92);
        box-shadow: 0 16px 40px rgba(15, 23, 42, 0.08);
      }}
      .label {{
        color: var(--accent);
        font-size: 12px;
        font-weight: 700;
        letter-spacing: 0.08em;
        text-transform: uppercase;
      }}
      section.slide h2 {{
        margin: 10px 0 18px;
        font-size: 28px;
      }}
      section.slide img {{
        display: block;
        width: 100%;
        border-radius: 16px;
        border: 1px solid var(--line);
      }}
      .text-slide {{
        border-radius: 18px;
        border: 1px dashed var(--line);
        padding: 18px 20px;
        background: #f8fafc;
      }}
      .text-slide .subtitle {{
        margin: 0 0 16px;
        color: var(--muted);
      }}
      .text-slide ul {{
        margin: 0;
        padding-left: 20px;
        line-height: 1.7;
      }}
      .text-slide .description {{
        margin: 16px 0 0;
        line-height: 1.7;
        color: var(--muted);
      }}
    </style>
  </head>
  <body>
    <main>
      <h1>{escape(str(plan.get("title") or output_path.stem))}</h1>
      <p class="meta">Auto-generated presentation preview</p>
      {''.join(sections)}
    </main>
  </body>
</html>
"""
    preview_path.write_text(html, encoding="utf-8")
    return preview_path


def generate_ppt(
    plan_file: str,
    slide_images: list[str] | None,
    output_file: str,
) -> str:
    """
    Generate a PowerPoint presentation from slide images.

    Args:
        plan_file: Path to JSON file containing presentation plan
        slide_images: List of paths to slide images in order
        output_file: Path to output PPTX file

    Returns:
        Status message
    """
    with open(plan_file, "r", encoding="utf-8") as f:
        plan = json.load(f)

    aspect_ratio = str(plan.get("aspect_ratio", "16:9"))
    slide_width, slide_height = _slide_dimensions(aspect_ratio)
    slide_images = slide_images or []

    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    slides_info = plan.get("slides", [])
    blank_layout = prs.slide_layouts[6]

    if not slides_info:
        slides_info = [{"title": Path(output_file).stem, "subtitle": "Generated presentation"}]

    normalized_images: list[str] = []
    for image_path in slide_images:
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"Slide image not found: {image_path}")
        normalized_images.append(image_path)

    for index, slide_info in enumerate(slides_info):
        slide = prs.slides.add_slide(blank_layout)
        image_path = normalized_images[index] if index < len(normalized_images) else None
        if image_path:
            _add_full_slide_image(slide, image_path, slide_width, slide_height)
        else:
            _add_text_slide(slide, slide_info, slide_width, slide_height)
        _add_notes(slide, slide_info)

    prs.save(output_file)
    preview_path = _render_preview_html(plan, normalized_images, output_file)

    fallback_count = max(len(slides_info) - len(normalized_images), 0)
    return (
        f"Successfully generated presentation with {len(slides_info)} slides to {output_file}. "
        f"HTML preview: {preview_path}. "
        f"Text fallback slides used: {fallback_count}."
    )


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Generate PowerPoint presentation from slide images"
    )
    parser.add_argument(
        "--plan-file",
        required=True,
        help="Absolute path to JSON presentation plan file",
    )
    parser.add_argument(
        "--slide-images",
        nargs="*",
        default=[],
        help="Absolute paths to slide images in order (optional, space-separated)",
    )
    parser.add_argument(
        "--output-file",
        required=True,
        help="Output path for generated PPTX file",
    )

    args = parser.parse_args()

    try:
        print(
            generate_ppt(
                args.plan_file,
                args.slide_images,
                args.output_file,
            )
        )
    except Exception as e:
        print(f"Error while generating presentation: {e}")
