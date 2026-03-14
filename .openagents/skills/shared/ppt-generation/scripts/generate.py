#!/usr/bin/env python3

import json
import os
from io import BytesIO
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

EMU_PER_INCH = 914400

THEMES = {
    "default": {
        "title": RGBColor(255, 255, 255),
        "subtitle": RGBColor(228, 232, 240),
        "body": RGBColor(245, 247, 250),
        "accent": RGBColor(95, 173, 255),
        "panel_fill": RGBColor(15, 23, 42),
        "panel_transparency": 18,
        "line": RGBColor(148, 163, 184),
    },
    "dark-premium": {
        "title": RGBColor(212, 175, 55),
        "subtitle": RGBColor(244, 228, 186),
        "body": RGBColor(255, 255, 255),
        "accent": RGBColor(212, 175, 55),
        "panel_fill": RGBColor(10, 10, 10),
        "panel_transparency": 12,
        "line": RGBColor(80, 62, 12),
    },
    "glassmorphism": {
        "title": RGBColor(255, 255, 255),
        "subtitle": RGBColor(230, 244, 255),
        "body": RGBColor(255, 255, 255),
        "accent": RGBColor(0, 212, 255),
        "panel_fill": RGBColor(255, 255, 255),
        "panel_transparency": 68,
        "line": RGBColor(255, 255, 255),
    },
    "gradient-modern": {
        "title": RGBColor(255, 255, 255),
        "subtitle": RGBColor(244, 244, 255),
        "body": RGBColor(255, 255, 255),
        "accent": RGBColor(255, 132, 91),
        "panel_fill": RGBColor(24, 24, 27),
        "panel_transparency": 24,
        "line": RGBColor(255, 132, 91),
    },
    "minimal-swiss": {
        "title": RGBColor(17, 24, 39),
        "subtitle": RGBColor(75, 85, 99),
        "body": RGBColor(31, 41, 55),
        "accent": RGBColor(220, 38, 38),
        "panel_fill": RGBColor(255, 255, 255),
        "panel_transparency": 8,
        "line": RGBColor(229, 231, 235),
    },
    "keynote": {
        "title": RGBColor(255, 255, 255),
        "subtitle": RGBColor(223, 230, 237),
        "body": RGBColor(255, 255, 255),
        "accent": RGBColor(66, 153, 225),
        "panel_fill": RGBColor(0, 0, 0),
        "panel_transparency": 24,
        "line": RGBColor(66, 153, 225),
    },
}


def resolve_slide_size(aspect_ratio: str) -> tuple[Any, Any]:
    if aspect_ratio == "4:3":
        return Inches(10), Inches(7.5)
    return Inches(13.333), Inches(7.5)


def resolve_theme(style: str | None) -> dict[str, Any]:
    theme = dict(THEMES["default"])
    if style and style in THEMES:
        theme.update(THEMES[style])
    return theme


def add_background_image(slide, image_path: str, slide_width, slide_height) -> None:
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Slide image not found: {image_path}")

    with Image.open(image_path) as img:
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")

        img_width, img_height = img.size
        img_aspect = img_width / img_height
        slide_aspect = int(slide_width) / int(slide_height)

        slide_width_emu = int(slide_width)
        slide_height_emu = int(slide_height)

        if img_aspect > slide_aspect:
            new_width_emu = slide_width_emu
            new_height_emu = int(slide_width_emu / img_aspect)
            left = Inches(0)
            top = Inches((slide_height_emu - new_height_emu) / EMU_PER_INCH)
        else:
            new_height_emu = slide_height_emu
            new_width_emu = int(slide_height_emu * img_aspect)
            left = Inches((slide_width_emu - new_width_emu) / EMU_PER_INCH)
            top = Inches(0)

        img_bytes = BytesIO()
        img.save(img_bytes, format="JPEG", quality=95)
        img_bytes.seek(0)

        slide.shapes.add_picture(
            img_bytes,
            left,
            top,
            Inches(new_width_emu / EMU_PER_INCH),
            Inches(new_height_emu / EMU_PER_INCH),
        )


def add_panel(slide, left: float, top: float, width: float, height: float, theme: dict[str, Any]):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = theme["panel_fill"]
    panel.fill.transparency = theme["panel_transparency"]
    panel.line.color.rgb = theme["line"]
    panel.line.width = Pt(1.25)
    return panel


def add_accent_bar(slide, left: float, top: float, width: float, theme: dict[str, Any]):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme["accent"]
    bar.line.fill.background()
    return bar


def add_textbox(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    font_color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)

    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    run.font.name = "Aptos"


def add_text_lines(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str],
    font_size: int,
    font_color: RGBColor,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)

    first = True
    for line in lines:
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(10)
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.name = "Aptos"


def slide_body_lines(slide_info: dict[str, Any]) -> list[str]:
    lines: list[str] = []

    for key in ("key_points", "bullet_points", "points"):
        value = slide_info.get(key)
        if isinstance(value, list):
            for item in value:
                text = stringify_text(item)
                if text:
                    lines.append(f"- {text}")

    for key in ("body", "content", "takeaway", "quote", "description"):
        value = slide_info.get(key)
        if isinstance(value, str):
            text = value.strip()
            if text:
                lines.append(text)
        elif isinstance(value, list):
            for item in value:
                text = stringify_text(item)
                if text:
                    lines.append(text)

    return lines


def stringify_text(value: Any) -> str | None:
    if isinstance(value, str):
        text = value.strip()
        return text or None
    if isinstance(value, dict):
        for key in ("text", "title", "content", "label"):
            text = value.get(key)
            if isinstance(text, str) and text.strip():
                return text.strip()
    return None


def write_notes(slide, slide_info: dict[str, Any]) -> None:
    notes: list[str] = []

    title = slide_info.get("title")
    subtitle = slide_info.get("subtitle")
    body_lines = slide_body_lines(slide_info)

    if isinstance(title, str) and title.strip():
        notes.append(f"Title: {title.strip()}")
    if isinstance(subtitle, str) and subtitle.strip():
        notes.append(f"Subtitle: {subtitle.strip()}")
    if body_lines:
        notes.append("Body:")
        notes.extend(f"  {line}" for line in body_lines)

    if not notes:
        return

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    if text_frame is not None:
        text_frame.text = "\n".join(notes)


def render_title_slide(slide, slide_info: dict[str, Any], theme: dict[str, Any]) -> None:
    add_panel(slide, 1.2, 1.45, 10.95, 4.5, theme)
    add_accent_bar(slide, 1.6, 1.85, 2.3, theme)

    title = str(slide_info.get("title") or "").strip()
    subtitle = str(slide_info.get("subtitle") or "").strip()

    if title:
        add_textbox(
            slide,
            left=1.7,
            top=2.15,
            width=10.15,
            height=1.8,
            text=title,
            font_size=28,
            font_color=theme["title"],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    if subtitle:
        add_textbox(
            slide,
            left=2.0,
            top=3.95,
            width=9.55,
            height=0.9,
            text=subtitle,
            font_size=16,
            font_color=theme["subtitle"],
            align=PP_ALIGN.CENTER,
        )


def render_content_slide(
    slide,
    slide_info: dict[str, Any],
    theme: dict[str, Any],
    slide_index: int,
) -> None:
    title = str(slide_info.get("title") or "").strip()
    subtitle = str(slide_info.get("subtitle") or "").strip()
    body_lines = slide_body_lines(slide_info)

    place_on_right = (slide_index + 1) % 2 == 0
    left = 7.05 if place_on_right else 0.95
    width = 5.3

    add_panel(slide, left, 0.8, width, 5.8, theme)
    add_accent_bar(slide, left + 0.35, 1.15, 1.5, theme)

    if title:
        add_textbox(
            slide,
            left=left + 0.35,
            top=1.35,
            width=width - 0.7,
            height=0.85,
            text=title,
            font_size=22,
            font_color=theme["title"],
            bold=True,
        )

    current_top = 2.2
    if subtitle:
        add_textbox(
            slide,
            left=left + 0.35,
            top=current_top,
            width=width - 0.7,
            height=0.6,
            text=subtitle,
            font_size=12,
            font_color=theme["subtitle"],
        )
        current_top += 0.65

    if body_lines:
        add_text_lines(
            slide,
            left=left + 0.4,
            top=current_top,
            width=width - 0.8,
            height=3.85,
            lines=body_lines,
            font_size=13,
            font_color=theme["body"],
        )


def render_slide(
    slide,
    slide_info: dict[str, Any],
    theme: dict[str, Any],
    slide_index: int,
) -> None:
    slide_type = str(slide_info.get("type") or "").strip().lower()
    body_lines = slide_body_lines(slide_info)
    subtitle = str(slide_info.get("subtitle") or "").strip()

    if slide_type == "title" or (slide_index == 0 and not body_lines):
        render_title_slide(slide, slide_info, theme)
        return

    if slide_type == "conclusion" and subtitle and not body_lines:
        render_title_slide(slide, slide_info, theme)
        return

    render_content_slide(slide, slide_info, theme, slide_index)


def generate_ppt(plan_file: str, slide_images: list[str], output_file: str) -> str:
    with open(plan_file, "r", encoding="utf-8") as f:
        plan = json.load(f)

    if not slide_images:
        raise ValueError("At least one slide image is required")

    slides_info = plan.get("slides", [])
    if slides_info and len(slide_images) != len(slides_info):
        raise ValueError(
            f"Slide image count {len(slide_images)} does not match plan slide count {len(slides_info)}"
        )

    slide_width, slide_height = resolve_slide_size(plan.get("aspect_ratio", "16:9"))
    theme = resolve_theme(plan.get("style"))

    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    blank_layout = prs.slide_layouts[6]

    for slide_index, image_path in enumerate(slide_images):
        slide = prs.slides.add_slide(blank_layout)
        add_background_image(slide, image_path, slide_width, slide_height)

        slide_info = slides_info[slide_index] if slide_index < len(slides_info) else {}
        if isinstance(slide_info, dict):
            render_slide(slide, slide_info, theme, slide_index)
            write_notes(slide, slide_info)

    prs.save(output_file)
    return (
        f"Successfully generated presentation with {len(slide_images)} slides and "
        f"native editable text to {output_file}"
    )


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Generate PowerPoint presentation from slide images"
    )
    parser.add_argument(
        "--plan-file",
        required=True,
        help="Absolute path to JSON presentation plan file",
    )
    parser.add_argument(
        "--slide-images",
        nargs="+",
        required=True,
        help="Absolute paths to slide images in order (space-separated)",
    )
    parser.add_argument(
        "--output-file",
        required=True,
        help="Output path for generated PPTX file",
    )

    args = parser.parse_args()

    try:
        print(
            generate_ppt(
                args.plan_file,
                args.slide_images,
                args.output_file,
            )
        )
    except Exception as e:
        print(f"Error while generating presentation: {e}")
